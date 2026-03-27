"""
METIS — Funzionalità PPT Generator
Screenshot di tutte le 16 funzionalità (M1-M8 + E1-E8) + Extra
"""

import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── PALETTE ──────────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0A, 0x0E, 0x27)
INDIGO     = RGBColor(0x1A, 0x1F, 0x4E)
ELECTRIC   = RGBColor(0x4F, 0x8E, 0xFF)
CYAN       = RGBColor(0x00, 0xD4, 0xE8)
GOLD       = RGBColor(0xFF, 0xC4, 0x30)
GREEN      = RGBColor(0x2E, 0xCC, 0x71)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xB0, 0xBA, 0xD0)
DARK_CARD  = RGBColor(0x12, 0x18, 0x3A)
PURPLE     = RGBColor(0xC0, 0x7A, 0xFF)

W = Inches(13.33)
H = Inches(7.5)

ARTIFACTS = r"C:\Users\GaetanoPecorella\.gemini\antigravity\brain\866479aa-49d0-47d7-bb91-4fca622c890f"

def img(name):
    """Find screenshot by partial name in artifacts dir."""
    for f in os.listdir(ARTIFACTS):
        if name in f and f.endswith('.png') and '.system' not in f:
            return os.path.join(ARTIFACTS, f)
    return None

def set_bg(slide, color):
    bg = slide.background; fill = bg.fill
    fill.solid(); fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); return s

def tb(slide, text, l, t, w, h, size=14, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    run.font.name = "Calibri"; return box

def header(slide, title, subtitle, accent=CYAN):
    rect(slide, 0, 0, W, Inches(1.05), INDIGO)
    rect(slide, 0, Inches(1.05), W, Pt(2.5), accent)
    tb(slide, title, Inches(0.35), Inches(0.05), Inches(9), Inches(0.6),
       size=26, bold=True, color=accent)
    if subtitle:
        tb(slide, subtitle, Inches(0.35), Inches(0.62), Inches(11), Inches(0.38),
           size=12, color=LIGHT_GREY)

def badge(slide, code, label, l, t, accent):
    r = rect(slide, l, t, Inches(0.65), Inches(0.32), accent)
    tb(slide, code, l, t-Pt(2), Inches(0.65), Inches(0.34),
       size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    tb(slide, label, l+Inches(0.72), t, Inches(3.5), Inches(0.34),
       size=12, bold=True, color=WHITE)

def slide_num(slide, n, tot):
    tb(slide, f"{n}/{tot}", Inches(12.5), Inches(7.15), Inches(0.8), Inches(0.28),
       size=10, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)

def add_screenshot(slide, path, l, t, w, h):
    if path and os.path.exists(path):
        slide.shapes.add_picture(path, l, t, w, h)
    else:
        r = rect(slide, l, t, w, h, DARK_CARD)
        tb(slide, "📸 Screenshot\nnon disponibile", l, t+h/2-Inches(0.3), w, Inches(0.6),
           size=13, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

def description_box(slide, lines, l, t, w, accent=ELECTRIC):
    """Multi-line description box with accent bar."""
    total_h = Inches(0.38 * len(lines)) + Inches(0.2)
    rect(slide, l, t, w, total_h, DARK_CARD)
    rect(slide, l, t, Pt(3), total_h, accent)
    box = slide.shapes.add_textbox(l+Inches(0.15), t+Inches(0.08), w-Inches(0.2), total_h-Inches(0.1))
    tf = box.text_frame; tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run(); run.text = line
        run.font.size = Pt(12); run.font.color.rgb = LIGHT_GREY
        run.font.name = "Calibri"

# ── SLIDE FACTORY ─────────────────────────────────────────────────────────────

def make_cover(prs, tot):
    """Cover slide con recap strutturato delle 3 categorie di funzionalità."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)

    # Top accent line
    rect(slide, 0, 0, W, Inches(0.05), ELECTRIC)

    # Header strip
    rect(slide, 0, Inches(0.05), W, Inches(1.0), INDIGO)
    tb(slide, "METIS  —  Panoramica Funzionalità",
       Inches(0.4), Inches(0.08), Inches(10), Inches(0.55),
       size=28, bold=True, color=CYAN)
    tb(slide, "Documento Finomnia AI in PEF V05 · 8 Moduli Core + 8 Evoluzioni + 3 Extra Metis  ·  Marzo 2026  ·  Confidenziale",
       Inches(0.4), Inches(0.63), Inches(12.5), Inches(0.34),
       size=11, color=LIGHT_GREY)
    rect(slide, 0, Inches(1.05), W, Pt(2), ELECTRIC)

    # ── COLONNA 1: 8 MODULI CORE ───────────────────────────────────────────
    col1_x = Inches(0.3)
    col_w  = Inches(4.18)
    rect(slide, col1_x, Inches(1.15), col_w, Inches(6.15), DARK_CARD)
    rect(slide, col1_x, Inches(1.15), col_w, Pt(3), ELECTRIC)

    # Header colonna
    rect(slide, col1_x, Inches(1.15), col_w, Inches(0.55), RGBColor(0x0D, 0x1A, 0x3E))
    tb(slide, "M1 – M8  ·  MODULI CORE", col1_x+Inches(0.12), Inches(1.19),
       col_w-Inches(0.15), Inches(0.4), size=12, bold=True, color=ELECTRIC)
    tb(slide, "Richiesti dal documento Finomnia AI in PEF V05",
       col1_x+Inches(0.12), Inches(1.57), col_w-Inches(0.15), Inches(0.28),
       size=9.5, italic=True, color=LIGHT_GREY)

    core_items = [
        ("M1", "Sintesi Commenti Storici PEF", "Analista del credito"),
        ("M2", "Nowcasting Reputazionale Web", "Analista del credito"),
        ("M3", "KPI Bilancio + Risk Score XGBoost", "Analista del credito"),
        ("M4", "Benchmark ATECO / ISTAT", "Analista / Resp. Crediti"),
        ("M5", "Analisi CR 12 Mesi (Anomaly Detection)", "Analista del credito"),
        ("M6", "Cross-Check CR ↔ Bilancio", "Analista / Compliance"),
        ("M7", "Forecast DSCR — 3 Scenari Prospettici", "Responsabile / Deliberante"),
        ("M8", "SWOT Matrix Auto-Generata", "Responsabile / Deliberante"),
    ]
    y0 = Inches(1.95)
    for code, name, user in core_items:
        # badge code
        rect(slide, col1_x+Inches(0.12), y0+Pt(2), Inches(0.48), Inches(0.29), ELECTRIC)
        tb(slide, code, col1_x+Inches(0.12), y0, Inches(0.48), Inches(0.31),
           size=9.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        tb(slide, name, col1_x+Inches(0.68), y0, col_w-Inches(0.8), Inches(0.31),
           size=11, bold=True, color=WHITE)
        tb(slide, user, col1_x+Inches(0.68), y0+Inches(0.28), col_w-Inches(0.8), Inches(0.26),
           size=9, color=LIGHT_GREY)
        y0 += Inches(0.64)

    # ── COLONNA 2: 8 EVOLUZIONI ────────────────────────────────────────────
    col2_x = Inches(4.57)
    rect(slide, col2_x, Inches(1.15), col_w, Inches(6.15), DARK_CARD)
    rect(slide, col2_x, Inches(1.15), col_w, Pt(3), CYAN)

    rect(slide, col2_x, Inches(1.15), col_w, Inches(0.55), RGBColor(0x06, 0x1F, 0x2A))
    tb(slide, "E1 – E8  ·  EVOLUZIONI", col2_x+Inches(0.12), Inches(1.19),
       col_w-Inches(0.15), Inches(0.4), size=12, bold=True, color=CYAN)
    tb(slide, "Implementate da Metis — roadmap Finomnia V05 pag. 11-12",
       col2_x+Inches(0.12), Inches(1.57), col_w-Inches(0.15), Inches(0.28),
       size=9.5, italic=True, color=LIGHT_GREY)

    evo_items = [
        ("E1", "Simulatore Prodotto Finanziario", "Gestore Commerciale"),
        ("E2", "Indicatori Pagamento per Cluster", "Analista Portafoglio"),
        ("E3", "Teoria dei Giochi — Nash Equilibrium", "Resp. Crediti Senior"),
        ("E4", "Early Warning AI (30/60/90 gg)", "Risk Manager"),
        ("E5", "Indicatori Rischio AR (Atipicità/Rilevanza)", "Compliance Officer"),
        ("E6", "Cluster Omogenei per Early Warning", "Risk Manager"),
        ("E7", "Lettura Allegati OCR + AI", "Analista del credito"),
        ("E8", "Verifica Aderenza Policy Creditizie", "Credit Officer"),
    ]
    y0 = Inches(1.95)
    for code, name, user in evo_items:
        rect(slide, col2_x+Inches(0.12), y0+Pt(2), Inches(0.48), Inches(0.29), CYAN)
        tb(slide, code, col2_x+Inches(0.12), y0, Inches(0.48), Inches(0.31),
           size=9.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        tb(slide, name, col2_x+Inches(0.68), y0, col_w-Inches(0.8), Inches(0.31),
           size=11, bold=True, color=WHITE)
        tb(slide, user, col2_x+Inches(0.68), y0+Inches(0.28), col_w-Inches(0.8), Inches(0.26),
           size=9, color=LIGHT_GREY)
        y0 += Inches(0.64)

    # ── COLONNA 3: EXTRA METIS ─────────────────────────────────────────────
    col3_x = Inches(8.84)
    col3_w = Inches(4.19)
    rect(slide, col3_x, Inches(1.15), col3_w, Inches(6.15), DARK_CARD)
    rect(slide, col3_x, Inches(1.15), col3_w, Pt(3), GOLD)

    rect(slide, col3_x, Inches(1.15), col3_w, Inches(0.55), RGBColor(0x1F, 0x18, 0x04))
    tb(slide, "EXTRA  ·  AGGIUNTE DA METIS", col3_x+Inches(0.12), Inches(1.19),
       col3_w-Inches(0.15), Inches(0.4), size=12, bold=True, color=GOLD)
    tb(slide, "Non richieste da Finomnia — aggiunte dalla vision di prodotto",
       col3_x+Inches(0.12), Inches(1.57), col3_w-Inches(0.15), Inches(0.28),
       size=9.5, italic=True, color=LIGHT_GREY)

    extra_items = [
        ("X1", "Visual Policy Builder",
         "Editor drag&drop no-code per configurare il flusso decisionale.\nBlocchi: Data Ingestion → Fraud → AI Scoring → Approve/Reject.\nWhat-If Simulator vs Champion policy.",
         "Resp. Crediti / IT banca", ELECTRIC),
        ("X2", "Audit Trail SHA-256 Immutabile",
         "Ogni azione genera un record con hash SHA-256.\nOverride analista obbligatoriamente motivato.\nEsportabile per Compliance Officer.",
         "Compliance Officer / Regolatore", CYAN),
        ("X3", "Fraud Detection su Documenti OCR",
         "ML Anomaly Detection su documenti (Zest Protect™-style).\nFalse Positive Limit configurabile (default 1.4%).\nIntegrato nel Visual Policy Builder.",
         "Analista / Anti-Frode", GREEN),
    ]
    y0 = Inches(1.95)
    for code, name, desc, user, ac in extra_items:
        rect(slide, col3_x+Inches(0.12), y0+Pt(2), Inches(0.48), Inches(0.29), ac)
        tb(slide, code, col3_x+Inches(0.12), y0, Inches(0.48), Inches(0.31),
           size=9.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        tb(slide, name, col3_x+Inches(0.68), y0, col3_w-Inches(0.82), Inches(0.31),
           size=11, bold=True, color=ac)
        tb(slide, desc, col3_x+Inches(0.12), y0+Inches(0.34), col3_w-Inches(0.2), Inches(0.72),
           size=9.5, color=LIGHT_GREY)
        tb(slide, f"👤 {user}", col3_x+Inches(0.12), y0+Inches(1.02), col3_w-Inches(0.2), Inches(0.26),
           size=9, bold=True, color=GOLD)
        y0 += Inches(1.72)

    # Bottom bar
    rect(slide, 0, Inches(7.2), W, Inches(0.3), INDIGO)
    tb(slide, "METIS v2.0  ·  Powered by Finomnia AI Research  ·  EU AI Act Compliant  ·  Totale: 19 funzionalità  ·  Marzo 2026",
       Inches(0.4), Inches(7.22), Inches(12.5), Inches(0.25),
       size=10, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

def make_func_slide(prs, n, tot, module_code, module_label, accent,
                    title, subtitle, screenshot_key,
                    what_it_does, how_it_works, user_role, output_list):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    header(slide, f"{module_code} — {title}", subtitle, accent)
    slide_num(slide, n, tot)

    # Badge
    badge(slide, module_code, module_label, Inches(11.8), Inches(0.12), accent)

    # Screenshot (left 55%)
    scr = img(screenshot_key)
    add_screenshot(slide, scr, Inches(0.35), Inches(1.18), Inches(7.3), Inches(4.65))

    # Right panel
    rx = Inches(7.85)
    rw = Inches(5.15)

    # Cosa fa
    tb(slide, "COSA FA", rx, Inches(1.18), rw, Inches(0.32),
       size=10, bold=True, color=accent)
    tb(slide, what_it_does, rx, Inches(1.5), rw, Inches(0.85),
       size=12, color=WHITE)

    # Come funziona
    tb(slide, "COME FUNZIONA", rx, Inches(2.42), rw, Inches(0.32),
       size=10, bold=True, color=accent)
    y_hw = Inches(2.74)
    for step in how_it_works:
        rect(slide, rx, y_hw+Pt(5), Pt(5), Pt(5), accent)
        tb(slide, step, rx+Inches(0.2), y_hw, rw-Inches(0.2), Inches(0.42),
           size=11, color=LIGHT_GREY)
        y_hw += Inches(0.44)

    # Output
    tb(slide, "OUTPUT", rx, y_hw+Inches(0.08), rw, Inches(0.32),
       size=10, bold=True, color=accent)
    y_out = y_hw+Inches(0.4)
    for out in output_list:
        rect(slide, rx+Inches(0.05), y_out+Pt(6), Pt(6), Pt(6), GREEN)
        tb(slide, out, rx+Inches(0.28), y_out, rw-Inches(0.3), Inches(0.38),
           size=11, color=WHITE)
        y_out += Inches(0.38)

    # User role footer
    rect(slide, Inches(0.35), Inches(6.05), Inches(7.3), Inches(0.45), INDIGO)
    rect(slide, Inches(7.85), Inches(6.05), Inches(5.15), Inches(0.45), INDIGO)
    tb(slide, f"👤  Utilizzatore: {user_role}",
       Inches(0.5), Inches(6.1), Inches(7.0), Inches(0.35),
       size=12, bold=True, color=GOLD)
    tb(slide, "EU AI Act · Human-in-the-Loop · XAI",
       Inches(8.0), Inches(6.1), Inches(4.8), Inches(0.35),
       size=11, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)

    # Bottom bar
    rect(slide, 0, Inches(6.55), W, Inches(0.95), DARK_CARD)
    rect(slide, 0, Inches(6.55), W, Pt(1.5), accent)

# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = W; prs.slide_height = H

    TOT = 17  # Cover + 2 contesto + M1-M8 + E1/E3/E4/E7 + X + Riepilogo

    # ─── SLIDE 1: Cover ─────────────────────────────────────────────────────
    print("Slide 1: Cover...")
    make_cover(prs, TOT)

    # ─── SLIDE 2-3: Contesto Generale ───────────────────────────────────────
    print("Slide 2: Command Center...")
    make_func_slide(prs, 2, TOT,
        "HOME", "Command Center", ELECTRIC,
        "Dashboard Portafoglio in Tempo Reale",
        "KPI aggregati, Portfolio Health, AI Recommendations live",
        "homepage_top",
        "Panoramica real-time dell'intero portafoglio crediti con KPI aggregati, distribuzione rischio e raccomandazioni AI prioritizzate.",
        ["KPI live: Avg PD 4.2%, Altman Z' 3.12, DSCR 1.45x, 142 pratiche attive",
         "Portfolio Health Score 71/100: Coverage Ratio 83%, Liquidity 91%",
         "AI Recommendations: alert prioritizzati ALTA / MEDIA / BASSA urgenza",
         "Pipeline Flow: Approvate / Pending / Rifiutate in tempo reale"],
        "Risk Manager / Responsabile Crediti / Direzione",
        ["Portfolio Health Score 0–100",
         "AI Recommendations con urgenza classificata",
         "KPI trend con delta vs periodo precedente"])

    print("Slide 3: Portafoglio...")
    make_func_slide(prs, 3, TOT,
        "PORTF.", "Portafoglio Creditizio", CYAN,
        "Gestione Portafoglio Aziende",
        "Lista pratiche con PD, Altman Z-Score, rischio e stato in tempo reale",
        "portafoglio_list",
        "Vista tabellare di tutte le aziende nel portafoglio crediti con score di rischio, distribuzione stati e accesso diretto alle pratiche.",
        ["Distribuzione stati: Approvata / In Analisi / Da Revisionare / Sospesa / Rifiutata",
         "Distribuzione rischio: Basso / Medio / Alto / Critico con treemap visiva",
         "PD Distribution con soglia 5.0% e range min–max",
         "Tabella ordinabile per PD, Altman Z, Rischio, Fatturato, Operatore"],
        "Analista del credito / Responsabile Crediti",
        ["Lista aziende con PD, Altman Z, stato pratica",
         "Accesso diretto alla Dashboard Analisi",
         "Export e filtri multi-criterio"])

    # ─── SLIDE 4-11: MODULI CORE M1–M8 ─────────────────────────────────────

    print("Slide 4: M1 Sintesi PEF...")
    make_func_slide(prs, 4, TOT,
        "M1", "Sintesi Commenti Storici PEF", ELECTRIC,
        "Sintesi Automatica dei Commenti PEF Storici",
        "LLM genera 3 paragrafi strutturati dalle ultime 2 revisioni PEF con anti-hallucination",
        "analysis_dashboard_full_top",
        "Legge i commenti degli analisti dalle ultime 2 revisioni PEF e genera una sintesi narrativa strutturata con evidenza delle variazioni significative.",
        ["Input: commenti testuali storici (bilancio + CR + andamentale)",
         "LLM Claude Sonnet / GPT-4o con prompt strutturato",
         "Anti-hallucination: ogni numero verificato contro il testo sorgente",
         "Source attribution: ogni paragrafo linkato alla frase originale"],
        "Analista del credito",
        ["sintesi_societaria: profilo in 3–5 righe",
         "sintesi_bilancio: andamento economico-patrimoniale",
         "delta_highlights: variazioni YoY con flag HIGH/MEDIUM/LOW",
         "confidence_scores: grado certezza per ogni sezione"])

    print("Slide 5: M2 Web Reputation...")
    make_func_slide(prs, 5, TOT,
        "M2", "Nowcasting Reputazionale Web", ELECTRIC,
        "Analisi Reputazione e Sentiment da Fonti Web",
        "Web scraping su Sole 24Ore, ANSA, Reuters + FinBERT sentiment finanziario italiano",
        "analysis_dashboard_bottom_e1_e7",
        "Cerca notizie e informazioni pubbliche sull'azienda su fonti whitelistate e produce un'analisi del sentiment reputazionale aggiornata al giorno dell'analisi.",
        ["Web scraping con rate limiting (2 req/sec) su fonti approvate",
         "FinBERT specializzato in linguaggio finanziario italiano",
         "Aggregazione sentiment pesata per autorevolezza fonte",
         "Identificazione automatica soggetti correlati dalle notizie"],
        "Analista del credito",
        ["Sentiment: POSITIVO / NEUTRO / NEGATIVO con score 0–1",
         "Lista articoli con URL, data, estratto e sentiment",
         "Soggetti correlati: persone o aziende emerse dalle notizie"])

    print("Slide 6: M3 KPI Bilancio...")
    make_func_slide(prs, 6, TOT,
        "M3", "KPI Bilancio + Risk Score", ELECTRIC,
        "Analisi KPI Finanziari con Risk Score XGBoost + SHAP",
        "9 KPI deterministici (Python puro, zero LLM) + Altman Z-Score + XGBoost + SHAP",
        "analysis_m3_bilancio_details",
        "Calcola deterministicamente 9 KPI finanziari confrontando due esercizi (T0 vs T-1), genera alert sulle soglie e produce uno score di rischio XGBoost con spiegazione SHAP.",
        ["9 KPI: EBITDA Margin, ROE, ROA, Leverage D/E, Current Ratio, Interest Coverage, NFP/EBITDA, DSO, DPO",
         "Trend per KPI: UP / STABLE / DOWN con variazione percentuale",
         "Risk Score 0–1 (XGBoost) con SHAP feature importance",
         "Altman Z-Score calcolato deterministicamente"],
        "Analista del credito",
        ["Tabella KPI con trend e alert soglie",
         "Risk Score 0–1 con SHAP waterfall chart",
         "Altman Z-Score: SAFE (>2.99) / GREY / DISTRESS"])

    print("Slide 7: M4 Benchmark ISTAT...")
    make_func_slide(prs, 7, TOT,
        "M4", "Benchmark ATECO / ISTAT", ELECTRIC,
        "Posizionamento Competitivo vs Mediane Settoriali ISTAT",
        "Database ISTAT ~200 codici ATECO — percentile esatto per ogni KPI",
        "analysis_dashboard_modules_part1",
        "Confronta i KPI dell'azienda con le mediane settoriali ISTAT per codice ATECO, determinando il posizionamento percentile nella distribuzione di settore.",
        ["Database ISTAT: mediane e percentili 25°/75° per ~200 codici ATECO",
         "Matching automatico: codice ATECO → gruppo settoriale",
         "Calcolo percentile esatto nella distribuzione settoriale",
         "Output alimenta automaticamente M7 (Forecast) e M8 (SWOT)"],
        "Analista del credito / Responsabile Crediti",
        ["Posizione per KPI: OUTPERFORMER / IN LINEA / UNDERPERFORMER",
         "Narrativa: 'EBITDA Margin 12% > mediana settoriale 8%'",
         "Credit Risk Radar: spider chart multi-dimensionale"])

    print("Slide 8: M5 CR 12 Mesi...")
    make_func_slide(prs, 8, TOT,
        "M5", "Analisi Centrale Rischi 12 Mesi", ELECTRIC,
        "Pattern Centrale Rischi — Anomaly Detection su 12 Mesi",
        "Z-score statistico su serie temporale CR — classifica anomalie TRANSITORIE vs STRUTTURALI",
        "analysis_dashboard_full_top",
        "Analizza 12 mesi di dati mensili CR (accordato, utilizzato, sconfinamenti, sofferenze, scaduto) identificando pattern anomali e comportamenti a rischio.",
        ["Time-series analysis: medie, delta, utilizzo medio su 12 mesi",
         "Anomaly detection Z-score: valori >2σ flaggati come anomali",
         "Classificazione: TRANSITORIA (rientrata <60gg) vs STRUTTURALE",
         "Utilizzo ratio: alert sopra 85% (rischio saturazione linee)"],
        "Analista del credito",
        ["Pattern CR 12 mesi con semaforo mensile (verde/giallo/rosso)",
         "Lista anomalie con Z-score, mese, valore atteso vs rilevato",
         "risk_score_cr: 0–1 con spiegazione XAI"])

    print("Slide 9: M6 Cross-Check...")
    make_func_slide(prs, 9, TOT,
        "M6", "Cross-Check CR ↔ Bilancio", ELECTRIC,
        "Rilevamento Incoerenze tra Bilancio e Centrale Rischi",
        "Delta% tra debiti bancari dichiarati a bilancio e dati reali CR — alert se >10%",
        "analysis_dashboard_full_top",
        "Confronta i debiti bancari dichiarati nel bilancio con i dati reali della Centrale Rischi, rilevando eventuali disallineamenti che potrebbero indicare omissioni o manipolazioni.",
        ["Calcolo delta %: |Bilancio - CR| / max(Bilancio, CR) × 100",
         "Soglie: 0–10% → OK | 10–25% → WARNING | >25% → CRITICAL",
         "Analisi cause automatica: sfasamento temporale, debiti non censiti",
         "Raccomandazione all'analista (es. 'Richiedere doc. integrativa')"],
        "Analista del credito / Compliance",
        ["delta_pct_totale: scarto percentuale bilancio vs CR",
         "Severity: OK / WARNING / CRITICAL con semaforo",
         "possibili_cause + raccomandazione azione"])

    print("Slide 10: M7 DSCR...")
    make_func_slide(prs, 10, TOT,
        "M7", "Forecast DSCR Prospettico", ELECTRIC,
        "Sostenibilità del Debito — 3 Scenari Prospettici 12 Mesi",
        "DSCR = (EBITDA - Tasse - Capex) / (Quota capitale + Interessi) su Ottimistico / Base / Stress",
        "analysis_dashboard_modules_part1",
        "Calcola il Debt Service Coverage Ratio prospettico su 3 scenari per valutare la capacità futura dell'azienda di servire il debito richiesto.",
        ["Formula: DSCR = (EBITDA proiettato - Tasse - Capex) / (Q.capitale + Interessi)",
         "Scenario selezionato automaticamente dagli output di M3 e M5",
         "Parametri configurabili: revenue growth, rate delta, cost growth",
         "Alert automatico se Stress DSCR < 1.0 (rischio insolvenza)"],
        "Responsabile Crediti / Deliberante",
        ["DSCR Ottimistico: 1.67x · Base: 1.45x · Stress: 1.02x",
         "Scenario auto-selezionato da AI con motivazione",
         "Probabilità di Default (PD) integrata: 2.1%"])

    print("Slide 11: M8 SWOT...")
    make_func_slide(prs, 11, TOT,
        "M8", "SWOT Matrix Auto-Generata", ELECTRIC,
        "Sintesi Strategica Automatica — Matrice SWOT con Evidenze Numeriche",
        "Aggrega M2+M3+M4+M5 in matrice 2×2 con item linkati al modulo sorgente",
        "dashboard_games_and_swot",
        "Aggrega automaticamente gli output di M2, M3, M4 e M5 in una matrice SWOT strutturata, con evidenza numerica per ogni punto e narrativa LLM executive.",
        ["FORZE: KPI sopra mediana ISTAT + CR regolare + sentiment positivo",
         "DEBOLEZZE: KPI in alert + anomalie CR strutturali + trend negativo",
         "OPPORTUNITÀ: trend settoriale positivo + macro favorevole",
         "MINACCE: sentiment web negativo + DSCR Stress <1.0"],
        "Responsabile Crediti / Deliberante",
        ["Matrice 2×2 con 3–5 item per quadrante + evidenza numerica",
         "source_module: ogni item linkato al modulo generante (M3, M4...)",
         "Narrativa LLM executive dei 4 quadranti"])

    # ─── SLIDE 12-15: EVOLUZIONI E1–E8 ─────────────────────────────────────

    print("Slide 12: E1 Simulatore...")
    make_func_slide(prs, 12, TOT,
        "E1", "Simulatore Prodotti PEF", CYAN,
        "Raccomandazione AI del Prodotto Creditizio Ottimale",
        "Confronto prodotti P-endo / P-uto / Mutuo / Leasing con score profittabilità e DSCR impatto",
        "dashboard_ccii_section",
        "Motore di raccomandazione AI che, analizzato il profilo creditizio, suggerisce il prodotto finanziario statisticamente più adatto e profittevole per la banca.",
        ["Analizza profilo cliente: settore, storico CR, DSCR, esposizioni",
         "Confronta prodotti: P-endo, P-uto, Mutuo, Leasing, Chirografario",
         "Calcola profittabilità attesa + DSCR impatto per ogni opzione",
         "Genera ranking con spiegazione XAI e score MEDIO/ALTO/CRITICO"],
        "Gestore Commerciale / Analista del credito",
        ["Prodotto raccomandato con punteggio e motivazione",
         "Confronto tabellare con Rata, Margine, DSCR impatto, TAEG",
         "Score rischio residuo per ogni alternativa"])

    print("Slide 13: E3 Teoria dei Giochi...")
    make_func_slide(prs, 13, TOT,
        "E3", "Teoria dei Giochi — Ottimizzazione Condizioni", CYAN,
        "Game Theory: Equilibrio di Nash per la Negoziazione del Credito",
        "Modello Stackelberg Banca-Debitore · Strategia ottimale su 3 scenari negoziali",
        "dashboard_games_and_swot",
        "Applica la teoria dei giochi al processo negoziale del credito per ottimizzare le condizioni offerte dalla banca in funzione delle possibili contromosse del cliente.",
        ["Modella negoziazione: gioco Banca vs Cliente con payoff definiti",
         "Calcola strategie dominanti e Equilibrio di Nash per profilo cliente",
         "3 strategie: Conservativa / Bilanciata / Aggressiva con payoff",
         "Simula scenari: Accetta Tutto / Negozia / Rifiuta"],
        "Responsabile Crediti Senior (pratiche >€1M)",
        ["Equilibrio di Nash evidenziato nella matrice payoff",
         "Struttura ottimale offerta: tasso, garanzie, durata, covenants",
         "Payoff Banca e Debitore per ogni combinazione strategica"])

    print("Slide 14: E4 Early Warning...")
    make_func_slide(prs, 14, TOT,
        "E4", "Early Warning AI — Monitoraggio Predittivo", CYAN,
        "Previsione Deterioramento Clienti a 30/60/90 Giorni",
        "XGBoost + LSTM su CR mensile + news web → score PD prospettico con alert automatico",
        "homepage_top",
        "Sistema di AI predittiva che monitora in continuo il portafoglio crediti e identifica precocemente i clienti a rischio di deterioramento nei prossimi 30–90 giorni.",
        ["Feed continuo: CR mensile, utilizzo linee, pagamenti, news web (M2)",
         "Modello ML (XGBoost + LSTM) addestrato su default storici",
         "Score probabilità di default a 30/60/90 giorni con SHAP",
         "Alert automatici all'analista responsabile con dashboard aggregata"],
        "Risk Manager / Responsabile Portafoglio",
        ["AI Recommendations: 3 livelli urgenza ALTA/MEDIA/BASSA",
         "Score PD prospettico con motivazione SHAP",
         "Dashboard portafoglio aggregata con early warning"])

    print("Slide 15: E7 Lettura Allegati...")
    make_func_slide(prs, 15, TOT,
        "E7", "Lettura Allegati Documentali — OCR+AI", CYAN,
        "Estrazione e Sintesi Automatica di PDF, Visure e Contratti Allegati",
        "LlamaParse/Tesseract OCR + AI synthesis → highlight clausole critiche",
        "dossier_upload",
        "Legge automaticamente i PDF allegati alla pratica (bilanci, visure, contratti, perizie) ed estrae gli elementi critici evidenziandoli all'analista.",
        ["OCR ad alta precisione su PDF scansionati (LlamaParse + Tesseract)",
         "Estrazione strutturata: date, importi, soggetti, clausole",
         "AI sintesi: individua clausole di accelerazione, covenant, garanzie",
         "Highlight automatico elementi che richiedono approfondimento"],
        "Analista del credito",
        ["Dati strutturati estratti da ogni documento allegato",
         "Sintesi AI con highlight clausole critiche",
         "Integrato nella pipeline OCR della colonna Raw Data Stream"])

    # ─── SLIDE 16: EXTRA X1–X3 ──────────────────────────────────────────────

    print("Slide 16: Rule Engine + Extra X...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    header(slide, "X1 X2 X3 — Funzionalità Extra Aggiunte da METIS",
           "Non richieste da Finomnia — sviluppate autonomamente dalla vision di prodotto", GOLD)
    slide_num(slide, 16, TOT)

    scr = img("rule_engine")
    add_screenshot(slide, scr, Inches(0.35), Inches(1.18), Inches(7.5), Inches(5.0))

    extras_def = [
        (ELECTRIC, "X1  🎛️  Visual Policy Builder",
         "Editor drag&drop no-code per configurare il flusso decisionale del credito.\nBlocchi: Data Ingestion → Fraud Detection → AI Scoring → Approve/Reject.\nWhat-If Simulator vs Champion policy. EU AI Act Fair Lending Score.\n👤  Responsabile Crediti / IT banca"),
        (CYAN, "X2  📜  Audit Trail SHA-256 Immutabile",
         "Ogni azione genera un record immutabile con hash SHA-256 di input e output.\nOverride analista obbligatoriamente motivato e loggato.\nEsportabile per Compliance Officer. Difendibile in ispezione regolamentare.\n👤  Compliance Officer / Regolatore"),
        (GOLD, "X3  🛡️  Fraud Detection su Documenti OCR",
         "ML Anomaly Detection su documenti (Zest Protect™-style).\nFalse Positive Limit configurabile (default 1.4%).\nIntegrato come blocco nel Visual Policy Builder.\n👤  Analista del credito / Anti-Frode"),
    ]
    y_e = Inches(1.25)
    for accent, title_e, desc_e in extras_def:
        rect(slide, Inches(8.05), y_e, Inches(5.0), Inches(1.62), DARK_CARD)
        rect(slide, Inches(8.05), y_e, Pt(3), Inches(1.62), accent)
        tb(slide, title_e, Inches(8.22), y_e+Inches(0.1), Inches(4.6), Inches(0.38),
           size=12, bold=True, color=accent)
        tb(slide, desc_e, Inches(8.22), y_e+Inches(0.48), Inches(4.6), Inches(1.1),
           size=10, color=LIGHT_GREY)
        y_e += Inches(1.72)

    # ─── SLIDE 17: Riepilogo ────────────────────────────────────────────────
    print("Slide 17: Riepilogo...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    rect(slide, 0, 0, W, Inches(0.06), ELECTRIC)
    rect(slide, 0, Inches(0.06), W, Inches(1.6), INDIGO)
    tb(slide, "Riepilogo Funzionalità METIS", Inches(0.5), Inches(0.15),
       Inches(11), Inches(0.6), size=28, bold=True, color=CYAN)
    tb(slide, "19 funzionalità implementate · 5 tipologie di utilizzatori · EU AI Act compliant",
       Inches(0.5), Inches(0.78), Inches(11), Inches(0.4), size=14, color=LIGHT_GREY)

    summary_data = [
        (ELECTRIC, "8 Moduli Core — M1–M8", [
            "M1  Sintesi PEF  —  Analista",
            "M2  Web Reputation  —  Analista",
            "M3  KPI + XGBoost + SHAP  —  Analista",
            "M4  Benchmark ISTAT  —  Analista / Resp.",
            "M5  CR 12 Mesi  —  Analista",
            "M6  Cross-Check CR  —  Analista / Compliance",
            "M7  DSCR Forecast  —  Deliberante",
            "M8  SWOT Matrix  —  Deliberante",
        ]),
        (CYAN, "8 Evoluzioni — E1–E8", [
            "E1  Simulatore Prodotti  —  Gestore",
            "E2  Cluster Pagamento  —  Analista",
            "E3  Teoria dei Giochi  —  Senior",
            "E4  Early Warning AI  —  Risk Manager",
            "E5  Rischio AR  —  Compliance",
            "E6  Cluster Omogenei  —  Risk Manager",
            "E7  OCR Allegati  —  Analista",
            "E8  Verifica Policy  —  Credit Officer",
        ]),
        (GOLD, "3 Extra — X1–X3 + KPI", [
            "X1  Visual Policy Builder  —  IT / Resp.",
            "X2  Audit Trail SHA-256  —  Compliance",
            "X3  Fraud Detection ML  —  Anti-Frode",
            "",
            "✅  EU AI Act nativo",
            "✅  Human-in-the-Loop",
            "✅  XAI su ogni output",
            "✅  < 30 secondi per pratica",
        ]),
    ]

    for i, (accent, group_title, items) in enumerate(summary_data):
        x = Inches(0.4 + i*4.3)
        rect(slide, x, Inches(1.85), Inches(4.15), Inches(5.4), DARK_CARD)
        rect(slide, x, Inches(1.85), Inches(4.15), Pt(3), accent)
        tb(slide, group_title, x+Inches(0.15), Inches(1.95), Inches(3.8), Inches(0.45),
           size=13, bold=True, color=accent)
        for j, item in enumerate(items):
            if item:
                rect(slide, x+Inches(0.2), Inches(2.5)+j*Inches(0.55)+Pt(7), Pt(5), Pt(5), accent)
                tb(slide, item, x+Inches(0.42), Inches(2.5)+j*Inches(0.55),
                   Inches(3.6), Inches(0.5), size=11, color=WHITE)

    rect(slide, 0, Inches(7.15), W, Inches(0.35), INDIGO)
    tb(slide, "METIS v2.0  ·  Powered by Finomnia AI Research  ·  EU AI Act Compliant  ·  Marzo 2026",
       Inches(0.5), Inches(7.18), Inches(12), Inches(0.28),
       size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    out = r"c:\Users\GaetanoPecorella\OneDrive - Altermaind\Desktop\METIS\Metis_Funzionalita_PPT.pptx"
    prs.save(out)
    print(f"\n✅ PPT salvato: {out}")
    return out

if __name__ == "__main__":
    main()
