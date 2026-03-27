"""
METIS — Investor Pitch Deck Generator (Slide 1–11)
"""

import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy
import os

# ── PALETTE ──────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0A, 0x0E, 0x27)   # background principale
INDIGO      = RGBColor(0x1A, 0x1F, 0x4E)   # sfondo secondario /sezioni
ELECTRIC    = RGBColor(0x4F, 0x8E, 0xFF)   # accent blu elettrico
CYAN        = RGBColor(0x00, 0xD4, 0xE8)   # accent cyan / highlight
GOLD        = RGBColor(0xFF, 0xC4, 0x30)   # accent caldo / warning
GREEN       = RGBColor(0x2E, 0xCC, 0x71)   # success / checkmark
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xB0, 0xBA, 0xD0)
DARK_CARD   = RGBColor(0x12, 0x18, 0x3A)   # card background

W = Inches(13.33)
H = Inches(7.5)

# ── HELPERS ──────────────────────────────────────────────────────────────────

def rgb_hex(rgb: RGBColor) -> str:
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"

def set_bg(slide, color: RGBColor):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_box(slide, text, l, t, w, h,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox

def add_para(tf, text, font_size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p

def slide_title_bar(slide, title, subtitle=None):
    """Top bar with title and optional subtitle."""
    add_rect(slide, 0, 0, W, Inches(1.1), INDIGO)
    add_text_box(slide, title,
                 Inches(0.4), Inches(0.08), Inches(10), Inches(0.6),
                 font_size=28, bold=True, color=CYAN)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.4), Inches(0.65), Inches(10), Inches(0.38),
                     font_size=13, color=LIGHT_GREY)

def add_slide_number(slide, num, total=11):
    add_text_box(slide, f"{num} / {total}",
                 Inches(12.4), Inches(7.1), Inches(0.8), Inches(0.3),
                 font_size=10, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)

def accent_line(slide, y=Inches(1.1)):
    """Thin horizontal accent line below title bar."""
    line = slide.shapes.add_shape(1, 0, y, W, Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = ELECTRIC
    line.line.fill.background()

def card(slide, l, t, w, h, fill=DARK_CARD):
    shape = add_rect(slide, l, t, w, h, fill)
    return shape

def multi_para_box(slide, l, t, w, h, items, title=None,
                   item_size=13, title_size=14, title_color=CYAN,
                   item_color=WHITE, bullet="•"):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = title_color
        run.font.name = "Calibri"
    for item in items:
        p = tf.add_paragraph() if (not first or title) else tf.paragraphs[0]
        if first and not title:
            first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(item_size)
        run.font.color.rgb = item_color
        run.font.name = "Calibri"
    return txBox

# ═══════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════

def make_slide_1_cover(prs):
    """COVER"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, NAVY)

    # gradient-like overlay top strip
    add_rect(slide, 0, 0, W, Inches(0.08), ELECTRIC)
    add_rect(slide, 0, Inches(0.08), W, Inches(2.5), INDIGO)

    # METIS wordmark
    add_text_box(slide, "METIS",
                 Inches(0.6), Inches(0.5), Inches(8), Inches(1.6),
                 font_size=90, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # accent dot
    dot = slide.shapes.add_shape(9, Inches(0.6), Inches(2.35), Inches(0.18), Inches(0.18))
    dot.fill.solid(); dot.fill.fore_color.rgb = CYAN
    dot.line.fill.background()

    add_text_box(slide, "AI Credit Intelligence for Italian Banks",
                 Inches(0.84), Inches(2.25), Inches(10), Inches(0.55),
                 font_size=22, bold=False, color=CYAN)

    add_text_box(slide,
                 "Il motore AI glass-box per il credit underwriting delle PMI italiane",
                 Inches(0.6), Inches(2.85), Inches(9), Inches(0.5),
                 font_size=16, color=LIGHT_GREY)

    # divider
    add_rect(slide, Inches(0.6), Inches(3.5), Inches(1.5), Pt(2), ELECTRIC)

    # mission strip
    add_text_box(slide,
                 "Da 6–9 ore a 30 secondi.  Spiegato.  Conforme.  Automatizzato.",
                 Inches(0.6), Inches(3.65), Inches(10), Inches(0.55),
                 font_size=17, bold=True, color=WHITE)

    # 3 pillars
    pillars = [
        ("🚀", "Velocità", "< 30 sec per pratica"),
        ("🔍", "Explainability", "XAI + SHAP su ogni output"),
        ("✅", "Compliance", "EU AI Act nativo"),
    ]
    for i, (icon, title, sub) in enumerate(pillars):
        x = Inches(0.6 + i * 4.2)
        card(slide, x, Inches(4.4), Inches(3.9), Inches(1.6))
        add_text_box(slide, icon + "  " + title,
                     x + Inches(0.15), Inches(4.5), Inches(3.6), Inches(0.45),
                     font_size=16, bold=True, color=CYAN)
        add_text_box(slide, sub,
                     x + Inches(0.15), Inches(4.9), Inches(3.6), Inches(0.45),
                     font_size=13, color=LIGHT_GREY)

    # bottom bar
    add_rect(slide, 0, Inches(6.85), W, Inches(0.65), INDIGO)
    add_text_box(slide, "Marzo 2026  ·  Confidenziale  ·  Altermaind",
                 Inches(0.5), Inches(6.9), Inches(8), Inches(0.45),
                 font_size=12, color=LIGHT_GREY)
    add_slide_number(slide, 1)


def make_slide_2_problema(prs):
    """IL PROBLEMA"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    slide_title_bar(slide, "Il Problema", "Il credito alle PMI italiane è ancora manuale, lento e opaco")
    accent_line(slide)
    add_slide_number(slide, 2)

    # 3 big number cards
    stats = [
        ("⏱️", "6–9 ore", "per analizzare una\npratica PMI in Italia"),
        ("📋", "500+", "uffici credito che usano\nancora Excel e PDF"),
        ("💶", "€180B", "impieghi PMI con analisi\nmanuale inefficiente"),
    ]
    for i, (icon, num, desc) in enumerate(stats):
        x = Inches(0.45 + i * 4.28)
        card(slide, x, Inches(1.25), Inches(3.98), Inches(2.1), DARK_CARD)
        add_text_box(slide, icon,
                     x + Inches(0.18), Inches(1.35), Inches(0.6), Inches(0.5), font_size=26)
        add_text_box(slide, num,
                     x + Inches(0.18), Inches(1.65), Inches(3.6), Inches(0.75),
                     font_size=38, bold=True, color=GOLD)
        add_text_box(slide, desc,
                     x + Inches(0.18), Inches(2.33), Inches(3.6), Inches(0.75),
                     font_size=12.5, color=LIGHT_GREY)

    # bullet points
    bullets = [
        "Gli analisti bancari dedicano l'80% del tempo alla raccolta dati manuale, non all'analisi",
        "Le pratiche PEF vengono costruite su bilanci vecchi di 18 mesi — nessun nowcasting",
        "Zero explainability: le banche non possono motivare ai clienti le decisioni AI (obbligo EU AI Act 2026)",
        "Il 73% del tempo lavorativo bancario è influenzabile dall'AI  (fonte: Accenture, 2025)",
        "Nessun player verticale italiano copre PEF + Centrale Rischi + XBRL in un unico sistema",
    ]
    add_rect(slide, Inches(0.45), Inches(3.55), W - Inches(0.9), Inches(3.55), INDIGO)
    y = Inches(3.7)
    for b in bullets:
        add_rect(slide, Inches(0.65), y + Pt(5), Pt(6), Pt(6), ELECTRIC)
        add_text_box(slide, b,
                     Inches(0.95), y, Inches(11.9), Inches(0.5),
                     font_size=14, color=WHITE)
        y += Inches(0.58)


def make_slide_3_soluzione(prs):
    """LA SOLUZIONE"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    slide_title_bar(slide, "La Soluzione", "METIS: da 6 ore a 30 secondi · Spiegato · Conforme · Automatizzato")
    accent_line(slide)
    add_slide_number(slide, 3)

    # Hero quote
    add_rect(slide, Inches(0.4), Inches(1.3), W - Inches(0.8), Inches(1.4), INDIGO)
    add_rect(slide, Inches(0.4), Inches(1.3), Pt(4), Inches(1.4), ELECTRIC)
    add_text_box(slide,
                 '"METIS è il primo motore AI glass-box verticale sul credito PMI italiano:\n'
                 '8 moduli AI integrati, EU AI Act compliant, human-in-the-loop by design."',
                 Inches(0.7), Inches(1.4), Inches(11.9), Inches(1.2),
                 font_size=17, italic=True, color=WHITE)

    # 3 pillars — large cards
    pillars = [
        ("🚀", "Velocità", ELECTRIC,
         ["Analisi completa in < 30 secondi",
          "vs 6–9 ore del processo manuale attuale",
          "Pipeline multi-agent asincrona"]),
        ("🔍", "Explainability", CYAN,
         ["SHAP feature importance su ogni score",
          "Source attribution: ogni paragrafo linkato alla fonte",
          "Confidence score per ogni sezione LLM"]),
        ("✅", "Compliance", GREEN,
         ["EU AI Act Regolamento UE 2024/1689 nativo",
          "Audit trail SHA-256 immutabile",
          "Human-in-the-loop obbligatorio"]),
    ]
    for i, (icon, title, accent, items) in enumerate(pillars):
        x = Inches(0.45 + i * 4.28)
        card(slide, x, Inches(2.9), Inches(3.98), Inches(3.8))
        add_rect(slide, x, Inches(2.9), Inches(3.98), Pt(4), accent)
        add_text_box(slide, icon + "  " + title,
                     x + Inches(0.2), Inches(3.0), Inches(3.6), Inches(0.55),
                     font_size=18, bold=True, color=accent)
        y2 = Inches(3.62)
        for item in items:
            add_rect(slide, x + Inches(0.25), y2 + Pt(6), Pt(5), Pt(5), accent)
            add_text_box(slide, item,
                         x + Inches(0.5), y2, Inches(3.3), Inches(0.52),
                         font_size=12, color=LIGHT_GREY)
            y2 += Inches(0.56)


def make_slide_4_prodotto(prs):
    """IL PRODOTTO — 8 MODULI"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    slide_title_bar(slide, "Il Prodotto: 8 Moduli ACE",
                    "Architettura modulare — pipeline AI integrata end-to-end")
    accent_line(slide)
    add_slide_number(slide, 4)

    modules = [
        ("M1", "Sintesi Commenti\nStorici PEF", ELECTRIC,
         "LLM + anti-hallucination\n+ confidence score"),
        ("M2", "Nowcasting\nReputazionale", CYAN,
         "Web scraping + FinBERT\n+ sentiment score"),
        ("M3", "KPI Bilancio +\nRisk Score", GOLD,
         "Deterministico Python\n+ XGBoost + SHAP"),
        ("M4", "Benchmark\nATECO / ISTAT", GREEN,
         "200 codici ATECO,\npercentile esatto"),
        ("M5", "Analisi CR\n12 Mesi", ELECTRIC,
         "Anomaly detection\n(Z-score statistico)"),
        ("M6", "Cross-Check\nCR vs Bilancio", CYAN,
         "Delta%, severity OK /\nWARNING / CRITICAL"),
        ("M7", "Forecast DSCR\nProspettico", GOLD,
         "3 scenari: Ottimistico\n/ Base / Stress"),
        ("M8", "SWOT Matrix\nAuto-Generata", GREEN,
         "Aggrega M2+M3+M4+M5,\nnarrativa LLM executive"),
    ]

    cols = 4
    cw = Inches(3.1)
    ch = Inches(1.55)
    gx = Inches(0.25)
    gy = Inches(0.2)
    start_x = Inches(0.5)
    start_y = Inches(1.3)

    for idx, (code, name, accent, desc) in enumerate(modules):
        row = idx // cols
        col = idx % cols
        x = start_x + col * (cw + gx)
        y = start_y + row * (ch + gy)
        card(slide, x, y, cw, ch)
        add_rect(slide, x, y, cw, Pt(3), accent)
        # Code badge
        add_rect(slide, x + Inches(0.12), y + Inches(0.12),
                 Inches(0.55), Inches(0.35), accent)
        add_text_box(slide, code,
                     x + Inches(0.12), y + Inches(0.1),
                     Inches(0.55), Inches(0.35),
                     font_size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text_box(slide, name,
                     x + Inches(0.75), y + Inches(0.08), Inches(2.2), Inches(0.6),
                     font_size=12.5, bold=True, color=WHITE)
        add_text_box(slide, desc,
                     x + Inches(0.12), y + Inches(0.78), cw - Inches(0.24), Inches(0.7),
                     font_size=11, color=LIGHT_GREY)

    # Bottom: 3 trasversali
    add_rect(slide, Inches(0.5), Inches(6.1), W - Inches(1.0), Pt(1.5), ELECTRIC)
    extras = [
        ("🧠  XAI", "SHAP + source attribution su ogni output"),
        ("📜  Audit Trail", "Record SHA-256 immutabile — Compliance ready"),
        ("🎛️  Visual Policy Builder", "configurazione no-code del flusso decisionale"),
    ]
    for i, (title, sub) in enumerate(extras):
        x = Inches(0.5 + i * 4.28)
        add_text_box(slide, title,
                     x, Inches(6.2), Inches(4.0), Inches(0.38),
                     font_size=13, bold=True, color=CYAN)
        add_text_box(slide, sub,
                     x, Inches(6.57), Inches(4.0), Inches(0.38),
                     font_size=11, color=LIGHT_GREY)


def make_slide_5_demo(prs):
    """DEMO & KEY OUTPUTS"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    slide_title_bar(slide, "La Pratica PEF che si Scrive da Sola",
                    "Ogni numero verificato · Ogni paragrafo linkato alla fonte")
    accent_line(slide)
    add_slide_number(slide, 5)

    # Left: mock dashboard
    add_rect(slide, Inches(0.4), Inches(1.25), Inches(6.5), Inches(5.85), DARK_CARD)
    add_rect(slide, Inches(0.4), Inches(1.25), Inches(6.5), Pt(3), ELECTRIC)
    add_text_box(slide, "🖥  METIS Dashboard — Pratica #2024-1187",
                 Inches(0.6), Inches(1.35), Inches(6.1), Inches(0.4),
                 font_size=12, bold=True, color=ELECTRIC)

    # Module status rows
    modules_status = [
        ("M1 Sintesi Commenti Storici", "✅ Completato", GREEN),
        ("M2 Nowcasting Reputazionale", "✅ Completato", GREEN),
        ("M3 KPI Bilancio + Risk Score", "✅ Completato", GREEN),
        ("M4 Benchmark ATECO", "✅ Completato", GREEN),
        ("M5 Analisi CR 12 Mesi", "⚠️ 1 anomalia WARNING", GOLD),
        ("M6 Cross-Check CR vs Bilancio", "✅ Delta 4.2% — OK", GREEN),
        ("M7 DSCR Forecast", "✅ Scenario Base: 1.32", GREEN),
        ("M8 SWOT Matrix", "✅ Generata", GREEN),
    ]
    for i, (name, status, color) in enumerate(modules_status):
        y = Inches(1.85) + i * Inches(0.54)
        add_text_box(slide, name, Inches(0.6), y, Inches(3.8), Inches(0.42),
                     font_size=11.5, color=LIGHT_GREY)
        add_text_box(slide, status, Inches(4.4), y, Inches(2.3), Inches(0.42),
                     font_size=11.5, color=color, align=PP_ALIGN.RIGHT)

    # Risk score bar mock
    add_rect(slide, Inches(0.6), Inches(6.45), Inches(5.8), Inches(0.2), INDIGO)
    add_rect(slide, Inches(0.6), Inches(6.45), Inches(1.6), Inches(0.2), GREEN)
    add_text_box(slide, "Risk Score: 0.28  (LOW RISK)",
                 Inches(0.6), Inches(6.68), Inches(5.5), Inches(0.35),
                 font_size=11, bold=True, color=GREEN)

    # Right: Key outputs
    outputs = [
        (ELECTRIC, "Risk Score 0–1 con SHAP Waterfall",
         "Feature importance: Leverage +0.23 · DSCR +0.18 · CR utilization +0.11"),
        (CYAN, "SWOT Auto-Generata con Evidenze",
         "Forza: EBITDA Margin 14% > mediana ISTAT 8% · Debolezza: CR utilization 88% (alert)"),
        (GOLD, "DSCR 3 Scenari (Ottimistico/Base/Stress)",
         "Ottimistico 1.65 · Base 1.32 · Stress 0.94 → flag Stress < 1.0"),
        (GREEN, "Report PDF bancario-standard",
         "Export PDF/HTML · EU AI Act disclaimer obbligatorio · Audit trail allegato"),
    ]
    x_r = Inches(7.2)
    y_r = Inches(1.3)
    for accent, title, desc in outputs:
        card(slide, x_r, y_r, Inches(5.75), Inches(1.35))
        add_rect(slide, x_r, y_r, Pt(3), Inches(1.35), accent)
        add_text_box(slide, title,
                     x_r + Inches(0.18), y_r + Inches(0.12), Inches(5.4), Inches(0.45),
                     font_size=13, bold=True, color=accent)
        add_text_box(slide, desc,
                     x_r + Inches(0.18), y_r + Inches(0.55), Inches(5.4), Inches(0.72),
                     font_size=11.5, color=LIGHT_GREY)
        y_r += Inches(1.47)


def make_slide_6_mercato(prs):
    """MERCATO"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    slide_title_bar(slide, "Il Mercato",
                    "Un mercato enorme, regolato e non ancora servito")
    accent_line(slide)
    add_slide_number(slide, 6)

    # TAM / SAM / SOM circles (represented as nested boxes)
    labels = [
        ("TAM", "€4.2B", "Credit AI globale\n(CAGR 24%, 2026)", ELECTRIC, Inches(0.4), Inches(1.25), Inches(6.2), Inches(5.8)),
        ("SAM", "€420M", "Europa meridionale\nbanche PMI", CYAN,    Inches(0.9), Inches(1.85), Inches(5.2), Inches(4.6)),
        ("SOM", "€42M",  "~500 istituti IT\nportafoglio PMI > €100M", GOLD, Inches(1.5), Inches(2.55), Inches(4.0), Inches(3.2)),
    ]
    for code, amount, desc, color, l, t, w, h in labels:
        shape = slide.shapes.add_shape(9, l, t, w, h)  # oval
        shape.fill.solid(); shape.fill.fore_color.rgb = DARK_CARD
        shape.line.color.rgb = color; shape.line.width = Pt(2)

    # Labels overlay TAM/SAM/SOM
    add_text_box(slide, "TAM  €4.2B",
                 Inches(0.6), Inches(1.38), Inches(3), Inches(0.5),
                 font_size=15, bold=True, color=ELECTRIC)
    add_text_box(slide, "Credit AI globale · CAGR 24%",
                 Inches(0.6), Inches(1.8), Inches(3), Inches(0.4),
                 font_size=10.5, color=LIGHT_GREY)

    add_text_box(slide, "SAM  €420M",
                 Inches(1.1), Inches(2.15), Inches(3), Inches(0.5),
                 font_size=14, bold=True, color=CYAN)
    add_text_box(slide, "Europa sud + IT · banche PMI",
                 Inches(1.1), Inches(2.55), Inches(3), Inches(0.4),
                 font_size=10.5, color=LIGHT_GREY)

    add_text_box(slide, "SOM  €42M",
                 Inches(1.8), Inches(3.05), Inches(2.6), Inches(0.5),
                 font_size=13, bold=True, color=GOLD)
    add_text_box(slide, "Target immediato:\n~500 istituti IT",
                 Inches(1.8), Inches(3.45), Inches(2.6), Inches(0.6),
                 font_size=10, color=LIGHT_GREY)

    # Right side: acceleratori
    add_text_box(slide, "Acceleratori Normativi",
                 Inches(7.0), Inches(1.3), Inches(6.0), Inches(0.5),
                 font_size=16, bold=True, color=CYAN)

    acceleratori = [
        (ELECTRIC, "EU AI Act 2026",
         "Obbligo explainability: i competitor black-box\nnon sono conformi → exit forzata dal mercato"),
        (CYAN, "Codice della Crisi d'Impresa",
         "Più obblighi di monitoraggio bancario proattivo\n→ bisogno di early warning automatizzati"),
        (GOLD, "Basilea IV",
         "Più capitale richiesto senza scoring accurato\n→ incentivo a investire in AI credit tools"),
    ]
    y_a = Inches(1.95)
    for accent, title, desc in acceleratori:
        card(slide, Inches(7.0), y_a, Inches(6.0), Inches(1.45))
        add_rect(slide, Inches(7.0), y_a, Pt(3), Inches(1.45), accent)
        add_text_box(slide, title,
                     Inches(7.18), y_a + Inches(0.12), Inches(5.6), Inches(0.4),
                     font_size=14, bold=True, color=accent)
        add_text_box(slide, desc,
                     Inches(7.18), y_a + Inches(0.5), Inches(5.6), Inches(0.85),
                     font_size=12, color=LIGHT_GREY)
        y_a += Inches(1.57)


def make_slide_7_competitor(prs):
    """POSIZIONAMENTO COMPETITIVO"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    slide_title_bar(slide, "Posizionamento Competitivo",
                    "METIS: il quadrante più strategico e meno affollato del mercato europeo")
    accent_line(slide)
    add_slide_number(slide, 7)

    # 2x2 map
    map_l, map_t = Inches(0.4), Inches(1.25)
    map_w, map_h = Inches(6.2), Inches(5.8)
    half_w = map_w / 2
    half_h = map_h / 2

    # quadrant backgrounds
    q_colors = [DARK_CARD, INDIGO, INDIGO, DARK_CARD]  # TL TR BL BR
    for qi, qc in enumerate(q_colors):
        row = qi // 2; col = qi % 2
        add_rect(slide,
                 map_l + col * half_w, map_t + row * half_h,
                 half_w, half_h, qc)

    # METIS quadrant highlight (top-right = high auto, high explainability)
    add_rect(slide, map_l + half_w, map_t, half_w, half_h,
             RGBColor(0x0A, 0x1F, 0x3A))
    add_rect(slide, map_l + half_w, map_t, half_w, Pt(3), ELECTRIC)
    add_rect(slide, map_l + half_w, map_t, Pt(3), half_h, ELECTRIC)
    add_text_box(slide, "✅ METIS",
                 map_l + half_w + Inches(0.3), map_t + Inches(0.25),
                 Inches(2.8), Inches(0.5),
                 font_size=18, bold=True, color=CYAN)
    add_text_box(slide, "Target Quadrant\nHigh Auto + High XAI",
                 map_l + half_w + Inches(0.3), map_t + Inches(0.7),
                 Inches(2.8), Inches(0.7),
                 font_size=11, color=LIGHT_GREY)

    # Other players
    add_text_box(slide, "Zest AI · Scienaptic\nACTICO",
                 map_l + Inches(0.15), map_t + Inches(0.3),
                 Inches(2.8), Inches(0.8),
                 font_size=12, color=LIGHT_GREY)
    add_text_box(slide, "FICO · Experian\nPowerCurve",
                 map_l + half_w + Inches(0.15), map_t + half_h + Inches(0.2),
                 Inches(2.8), Inches(0.8),
                 font_size=12, color=LIGHT_GREY)
    add_text_box(slide, "Processi bancari\ntradizionali",
                 map_l + Inches(0.15), map_t + half_h + Inches(0.3),
                 Inches(2.8), Inches(0.7),
                 font_size=12, color=LIGHT_GREY)

    # Axis labels
    add_text_box(slide, "← Bassa Explainability",
                 map_l, map_t + map_h + Inches(0.05), Inches(3.0), Inches(0.38),
                 font_size=10.5, color=LIGHT_GREY)
    add_text_box(slide, "Alta Explainability →",
                 map_l + half_w, map_t + map_h + Inches(0.05), Inches(3.0), Inches(0.38),
                 font_size=10.5, color=ELECTRIC)
    add_text_box(slide, "↑ Alta Automazione",
                 map_l + map_w + Inches(0.08), map_t, Inches(1.8), Inches(0.38),
                 font_size=10, color=LIGHT_GREY)
    add_text_box(slide, "↓ Bassa Automazione",
                 map_l + map_w + Inches(0.08), map_t + half_h, Inches(1.8), Inches(0.38),
                 font_size=10, color=LIGHT_GREY)

    # Comparison table
    tbl_l = Inches(7.1)
    headers = ["Feature", "METIS", "FICO", "Zest AI", "Tradicionale"]
    col_widths = [Inches(2.3), Inches(0.8), Inches(0.8), Inches(0.8), Inches(1.2)]
    rows_data = [
        ("XAI / Glass-box",        "✅", "⚠️", "❌", "❌"),
        ("Verticale IT (PEF/CR)",  "✅", "❌", "❌", "❌"),
        ("EU AI Act compliant",    "✅", "⚠️", "❌", "N/A"),
        ("Policy Builder no-code", "✅", "❌", "❌", "❌"),
        ("< 30 sec per pratica",   "✅", "❌", "⚠️", "❌"),
    ]

    # Header row
    x = tbl_l
    for j, (hdr, cw) in enumerate(zip(headers, col_widths)):
        color = CYAN if j == 0 else (GREEN if j == 1 else LIGHT_GREY)
        add_rect(slide, x, Inches(1.3), cw, Inches(0.42), INDIGO)
        add_text_box(slide, hdr, x + Inches(0.05), Inches(1.33), cw, Inches(0.38),
                     font_size=11, bold=True, color=color, align=(PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT))
        x += cw

    # Data rows
    for i, (feat, *vals) in enumerate(rows_data):
        bg = DARK_CARD if i % 2 == 0 else INDIGO
        x = tbl_l
        y_row = Inches(1.72) + i * Inches(0.52)
        add_rect(slide, x, y_row, sum(col_widths), Inches(0.5), bg)
        add_text_box(slide, feat, x + Inches(0.05), y_row + Inches(0.06), col_widths[0], Inches(0.42),
                     font_size=11.5, color=WHITE)
        x += col_widths[0]
        for j, (val, cw) in enumerate(zip(vals, col_widths[1:])):
            vcol = GREEN if val == "✅" else (GOLD if val == "⚠️" else LIGHT_GREY)
            add_text_box(slide, val, x + Inches(0.05), y_row + Inches(0.06), cw, Inches(0.42),
                         font_size=14, bold=True, color=vcol, align=PP_ALIGN.CENTER)
            x += cw


def make_slide_8_usp(prs):
    """USP & MOAT"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    slide_title_bar(slide, "USP & Vantaggio Competitivo",
                    "5 vantaggi strutturali difficili da replicare")
    accent_line(slide)
    add_slide_number(slide, 8)

    usps = [
        ("🔍", "Glass-Box First", ELECTRIC,
         "XAI su ogni output — non un add-on ma architettura core. "
         "SHAP waterfall + source attribution su ogni paragrafo LLM."),
        ("🇮🇹", "Vertical Italian Market", CYAN,
         "PEF, Centrale Rischi, XBRL, ISTAT ATECO integrati nativamente. "
         "Nessun competitor straniero conosce le specificità del credito bancario italiano."),
        ("✅", "EU AI Act First Mover", GREEN,
         "Compliance nativa, non retrofittata. Dal 2026 sarà obbligo di legge: "
         "i competitor black-box escono dal mercato, METIS cattura la domanda."),
        ("📊", "Data Moat Proprietario", GOLD,
         "Dataset storici PEF + benchmark ATECO proprietari che crescono con ogni banca onboardata. "
         "Più clienti = modelli più accurati = vantaggio competitivo cumulativo."),
        ("🎛️", "Visual Policy Builder", RGBColor(0xC0, 0x7A, 0xFF),
         "Editor drag & drop no-code per il flusso decisionale. "
         "IP depositabile. Unico nel segmento bancario italiano."),
    ]

    for i, (icon, title, accent, desc) in enumerate(usps):
        y = Inches(1.35) + i * Inches(1.19)
        add_rect(slide, Inches(0.4), y, W - Inches(0.8), Inches(1.1), DARK_CARD)
        add_rect(slide, Inches(0.4), y, Pt(4), Inches(1.1), accent)
        # number badge
        badge = slide.shapes.add_shape(9, Inches(0.7), y + Inches(0.28),
                                        Inches(0.52), Inches(0.52))
        badge.fill.solid(); badge.fill.fore_color.rgb = accent
        badge.line.fill.background()

        add_text_box(slide, str(i+1),
                     Inches(0.7), y + Inches(0.28), Inches(0.52), Inches(0.52),
                     font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

        add_text_box(slide, icon + "  " + title,
                     Inches(1.35), y + Inches(0.08), Inches(3.8), Inches(0.45),
                     font_size=15, bold=True, color=accent)
        add_text_box(slide, desc,
                     Inches(1.35), y + Inches(0.52), Inches(11.3), Inches(0.55),
                     font_size=12.5, color=LIGHT_GREY)

    # Quote
    add_rect(slide, Inches(0.4), Inches(7.08), W - Inches(0.8), Inches(0.3), INDIGO)
    add_text_box(slide,
                 '"Con l\'EU AI Act, i competitor black-box andranno incontro a frizioni normative significative. '
                 'METIS è strutturalmente posizionato per catturare questa transizione."',
                 Inches(0.6), Inches(7.1), Inches(12.0), Inches(0.28),
                 font_size=11, italic=True, color=LIGHT_GREY)


def make_slide_9_business_model(prs):
    """BUSINESS MODEL"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    slide_title_bar(slide, "Business Model",
                    "SaaS B2B con revenue ricorrente e upsell per volumi")
    accent_line(slide)
    add_slide_number(slide, 9)

    # 3 pricing tiers
    tiers = [
        ("Starter", ELECTRIC, "Licenza base\nfino a 500 pratiche/anno",
         "€30K – 50K / anno", ["Tutti gli 8 moduli ACE", "Dashboard + report PDF", "SLA 99.5%"]),
        ("Growth ★", CYAN, "Per pratica elaborata\nvia API (pay-per-use)",
         "€10 – 25 / pratica", ["API REST nativa", "Volume discount", "Ideal: 1.000–5.000 prat/anno"]),
        ("Enterprise", GOLD, "Licensing + custom\nimplementazione + training",
         "€150K+ / anno", ["Connettori CBS custom", "On-premise option", "Dedicated CSM"]),
    ]
    for i, (name, accent, desc, price, features) in enumerate(tiers):
        x = Inches(0.45 + i * 4.28)
        card(slide, x, Inches(1.3), Inches(3.98), Inches(4.1), DARK_CARD)
        add_rect(slide, x, Inches(1.3), Inches(3.98), Pt(4), accent)
        add_text_box(slide, name,
                     x + Inches(0.2), Inches(1.4), Inches(3.6), Inches(0.5),
                     font_size=18, bold=True, color=accent)
        add_text_box(slide, desc,
                     x + Inches(0.2), Inches(1.88), Inches(3.6), Inches(0.65),
                     font_size=12, color=LIGHT_GREY)
        add_rect(slide, x + Inches(0.2), Inches(2.6), Inches(3.55), Pt(1), accent)
        add_text_box(slide, price,
                     x + Inches(0.2), Inches(2.68), Inches(3.6), Inches(0.55),
                     font_size=16, bold=True, color=WHITE)
        y2 = Inches(3.3)
        for feat in features:
            add_rect(slide, x + Inches(0.25), y2 + Pt(6), Pt(5), Pt(5), accent)
            add_text_box(slide, feat,
                         x + Inches(0.48), y2, Inches(3.3), Inches(0.45),
                         font_size=12, color=LIGHT_GREY)
            y2 += Inches(0.5)

    # GTM + unit economics strip
    add_rect(slide, Inches(0.4), Inches(5.65), W - Inches(0.8), Inches(1.65), INDIGO)
    add_text_box(slide, "Go-To-Market",
                 Inches(0.6), Inches(5.75), Inches(6), Inches(0.4),
                 font_size=13, bold=True, color=CYAN)
    gtm = [
        "Partnership CBS primari (Temenos, CEDACRI, CSE) come canale primario",
        "Pilota gratuito 60 gg con 2–3 BCC → conversion → scale",
        "Year 1: 5–8 pilot  ·  Year 2: 30+ contratti attivi",
    ]
    for j, g in enumerate(gtm):
        add_text_box(slide, "›  " + g,
                     Inches(0.6), Inches(6.18) + j * Inches(0.34), Inches(7.5), Inches(0.33),
                     font_size=11.5, color=WHITE)

    add_text_box(slide, "Unit Economics",
                 Inches(8.2), Inches(5.75), Inches(4.8), Inches(0.4),
                 font_size=13, bold=True, color=GOLD)
    ue = [
        "ARR per banca (Growth): ~€20K–50K",
        "Gross Margin target: 75–80%",
        "Payback period stimato: 9–14 mesi",
    ]
    for j, u in enumerate(ue):
        add_text_box(slide, "›  " + u,
                     Inches(8.2), Inches(6.18) + j * Inches(0.34), Inches(4.8), Inches(0.33),
                     font_size=11.5, color=WHITE)


def make_slide_10_trazione(prs):
    """TRAZIONE & VALIDAZIONE"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    slide_title_bar(slide, "Trazione & Validazione",
                    "Già costruito. Già testato. In validazione con il mercato.")
    accent_line(slide)
    add_slide_number(slide, 10)

    # Left: milestones
    add_text_box(slide, "Milestone Raggiunte",
                 Inches(0.4), Inches(1.3), Inches(6.5), Inches(0.45),
                 font_size=15, bold=True, color=CYAN)
    milestones = [
        (GREEN, "8 moduli ACE implementati e funzionanti"),
        (GREEN, "Anti-hallucination validation su ogni output LLM"),
        (GREEN, "Gap Analysis vs specifica Finomnia: Metis supera ogni modulo richiesto"),
        (GREEN, "Compliance EU AI Act documentata e audit-ready"),
        (GREEN, "Dashboard UI live con Visual Policy Builder integrato"),
        (GREEN, "Audit Trail SHA-256 immutabile operativo"),
        (GREEN, "Data Quality Scoring automatico su ogni fonte"),
        (GREEN, "Report PDF export in formato bancario standard"),
    ]
    for i, (color, text) in enumerate(milestones):
        y = Inches(1.85) + i * Inches(0.6)
        badge = slide.shapes.add_shape(9, Inches(0.55), y + Pt(5),
                                        Inches(0.35), Inches(0.35))
        badge.fill.solid(); badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        add_text_box(slide, "✓",
                     Inches(0.55), y + Pt(3), Inches(0.35), Inches(0.38),
                     font_size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text_box(slide, text,
                     Inches(1.05), y, Inches(5.7), Inches(0.5),
                     font_size=12.5, color=WHITE)

    # Right: metrics / pipeline
    add_text_box(slide, "Validazione Tecnica",
                 Inches(7.1), Inches(1.3), Inches(5.9), Inches(0.45),
                 font_size=15, bold=True, color=GOLD)

    metrics = [
        ("< 30 sec", "Tempo medio per pratica completa\n(vs 6–9 ore manuale)"),
        ("8 / 8", "Moduli Finomnia V05 pienamente\ncoperti e superati da Metis ACE"),
        ("0", "Allucinazioni non rilevate su test set\ncon validation anti-hallucination attiva"),
        ("SHA-256", "Standard di hashing audit trail\nimmutabile per ogni azione"),
    ]
    y_m = Inches(1.85)
    for val, desc in metrics:
        card(slide, Inches(7.1), y_m, Inches(5.9), Inches(1.2), DARK_CARD)
        add_text_box(slide, val,
                     Inches(7.25), y_m + Inches(0.1), Inches(2.5), Inches(0.6),
                     font_size=28, bold=True, color=GOLD)
        add_text_box(slide, desc,
                     Inches(9.3), y_m + Inches(0.14), Inches(3.55), Inches(0.9),
                     font_size=11.5, color=LIGHT_GREY)
        y_m += Inches(1.3)

    # Pipeline placeholder
    add_rect(slide, Inches(0.4), Inches(6.7), W - Inches(0.8), Inches(0.65), INDIGO)
    add_text_box(slide, "⚡  Pipeline Commerciale: in discussione con istituti bancari italiani — dettagli disponibili su richiesta",
                 Inches(0.6), Inches(6.79), Inches(12.5), Inches(0.45),
                 font_size=12, italic=True, color=LIGHT_GREY)


def make_slide_11_roadmap(prs):
    """ROADMAP"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    slide_title_bar(slide, "Roadmap 2026–2027",
                    "Espansione modulare e internazionalizzazione")
    accent_line(slide)
    add_slide_number(slide, 11)

    # Timeline bar
    phases = [
        ("Q2 2026", ELECTRIC, [
            "Onboarding 3 BCC pilot",
            "Pilota gratuito 60 gg",
            "Feedback loop analisti",
        ]),
        ("Q3 2026", CYAN, [
            "Connettori nativi:\nXBRL, Bankitalia API,\nCerved API",
            "Monitoraggio continuo\n12 mesi post-erogazione",
            "Early Warning System",
        ]),
        ("Q4 2026", GOLD, [
            "ESG Integration scoring\n(conformità EBA)",
            "Connettore CBS nativo\nTemenos / CEDACRI",
            "30+ contratti attivi",
        ]),
        ("2027", GREEN, [
            "Espansione Europa Sud:\nSpagna, Portogallo",
            "Open Banking API\nPSD2 integration",
            "Series A / Scale-up",
        ]),
    ]

    # Horizontal line
    line_y = Inches(3.3)
    add_rect(slide, Inches(0.5), line_y, W - Inches(1.0), Pt(3), ELECTRIC)

    for i, (label, accent, items) in enumerate(phases):
        x = Inches(0.5 + i * 3.2)
        # dot on timeline
        dot = slide.shapes.add_shape(9, x + Inches(1.35), line_y - Inches(0.18),
                                      Inches(0.36), Inches(0.36))
        dot.fill.solid(); dot.fill.fore_color.rgb = accent
        dot.line.fill.background()

        # label above line
        add_text_box(slide, label,
                     x, line_y - Inches(0.85), Inches(3.0), Inches(0.55),
                     font_size=17, bold=True, color=accent, align=PP_ALIGN.CENTER)

        # card below line
        card(slide, x, line_y + Inches(0.3), Inches(3.05), Inches(3.7), DARK_CARD)
        add_rect(slide, x, line_y + Inches(0.3), Inches(3.05), Pt(3), accent)
        y2 = line_y + Inches(0.45)
        for item in items:
            add_rect(slide, x + Inches(0.18), y2 + Pt(7), Pt(5), Pt(5), accent)
            add_text_box(slide, item,
                         x + Inches(0.38), y2, Inches(2.55), Inches(0.78),
                         font_size=11.5, color=WHITE)
            y2 += Inches(0.9)

    # Bottom: priority evolutions
    add_rect(slide, Inches(0.4), Inches(7.08), W - Inches(0.8), Inches(0.33), INDIGO)
    add_text_box(slide,
                 "Priorità roadmap: Monitoring Continuo · Early Warning · ESG Scoring · Connettori CBS nativi",
                 Inches(0.6), Inches(7.1), Inches(12.0), Inches(0.28),
                 font_size=11.5, bold=True, color=CYAN)


# ═══════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    print("Generando slide 1: Cover...")
    make_slide_1_cover(prs)
    print("Generando slide 2: Il Problema...")
    make_slide_2_problema(prs)
    print("Generando slide 3: La Soluzione...")
    make_slide_3_soluzione(prs)
    print("Generando slide 4: Il Prodotto...")
    make_slide_4_prodotto(prs)
    print("Generando slide 5: Demo...")
    make_slide_5_demo(prs)
    print("Generando slide 6: Mercato...")
    make_slide_6_mercato(prs)
    print("Generando slide 7: Competitor...")
    make_slide_7_competitor(prs)
    print("Generando slide 8: USP & Moat...")
    make_slide_8_usp(prs)
    print("Generando slide 9: Business Model...")
    make_slide_9_business_model(prs)
    print("Generando slide 10: Trazione...")
    make_slide_10_trazione(prs)
    print("Generando slide 11: Roadmap...")
    make_slide_11_roadmap(prs)

    out = r"c:\Users\GaetanoPecorella\OneDrive - Altermaind\Desktop\METIS\Metis_InvestorDeck_2026.pptx"
    prs.save(out)
    print(f"\n✅ Pitch deck salvato: {out}")

if __name__ == "__main__":
    main()
