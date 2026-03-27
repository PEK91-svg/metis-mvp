from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import os

BRAIN = r"C:\Users\GaetanoPecorella\.gemini\antigravity\brain\9473e7ce-4045-485e-a75e-2c69e4299e63"

IMAGES = {
    "e1": os.path.join(BRAIN, "e1_simulatore_prodotto_1774533275255.png"),
    "e2": os.path.join(BRAIN, "e2_cluster_pagamento_1774533296592.png"),
    "e3": os.path.join(BRAIN, "e3_teoria_giochi_1774533363088.png"),
    "e4": os.path.join(BRAIN, "e4_early_warning_1774533315539.png"),
    "e5": os.path.join(BRAIN, "e5_rischio_ar_1774533382708.png"),
    "e6": os.path.join(BRAIN, "e6_cluster_clientela_1774533332020.png"),
    "e7": os.path.join(BRAIN, "e7_ocr_allegati_1774533400226.png"),
    "e8": os.path.join(BRAIN, "e8_aderenza_policy_1774533420226.png"),
}

# Colors
BG_DARK   = RGBColor(0x09, 0x0E, 0x1C)   # dark navy
CYAN      = RGBColor(0x00, 0xD4, 0xFF)    # neon cyan
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xB0, 0xC8, 0xE0)
GOLD      = RGBColor(0xFF, 0xC0, 0x00)
GREEN     = RGBColor(0x4C, 0xE0, 0x8A)
ORANGE    = RGBColor(0xFF, 0x7A, 0x00)
GRAY      = RGBColor(0x6B, 0x7B, 0x8D)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

# ── Layout helper ────────────────────────────────────────────────
blank_layout = prs.slide_layouts[6]  # completely blank

def set_bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, text, left, top, width, height,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, font="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    from pptx.util import Pt
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = font
    return txBox

def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_line(slide, x1, y1, x2, y2, color=CYAN, width_pt=1.5):
    from pptx.util import Pt
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)

# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)

# Top accent bar
add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), CYAN)

# Left gradient panel
add_rect(slide, 0, 0, Inches(4.5), SLIDE_H, RGBColor(0x05, 0x18, 0x35))

# Logo text
add_text(slide, "METIS", Inches(0.4), Inches(1.0), Inches(3.5), Inches(1.2),
         size=54, bold=True, color=CYAN, font="Calibri")
add_text(slide, "ACE — Analisi Complementare Evoluta", Inches(0.4), Inches(2.0), Inches(3.6), Inches(0.8),
         size=13, bold=False, color=LIGHT, font="Calibri")

# Divider
add_rect(slide, Inches(0.4), Inches(2.9), Inches(2.0), Inches(0.04), CYAN)

# Subtitle
add_text(slide, "Roadmap Evolutiva", Inches(0.4), Inches(3.1), Inches(3.8), Inches(0.8),
         size=22, bold=True, color=WHITE, font="Calibri")
add_text(slide, "8 Evoluzioni Future del Prodotto\ndal documento Finomnia AI in PEF V05",
         Inches(0.4), Inches(3.75), Inches(3.8), Inches(1.0),
         size=13, color=LIGHT, font="Calibri")

add_text(slide, "Marzo 2026  |  Confidenziale", Inches(0.4), Inches(6.8), Inches(3.5), Inches(0.4),
         size=10, color=GRAY, font="Calibri")

# Right side — 8 pill boxes
pills = [
    ("E1", "Simulatore Prodotto"),
    ("E2", "Cluster Pagamento"),
    ("E3", "Teoria dei Giochi"),
    ("E4", "Early Warning AI"),
    ("E5", "Rischio AR"),
    ("E6", "Cluster Clientela"),
    ("E7", "OCR Allegati"),
    ("E8", "Verifica Policy"),
]
colors_pill = [CYAN, RGBColor(0xA0,0x5C,0xFF), GOLD, GREEN,
               ORANGE, CYAN, GREEN, RGBColor(0xFF,0x50,0x80)]

for i, (code, label) in enumerate(pills):
    row = i % 4; col = i // 4
    lft = Inches(4.9 + col * 4.1)
    tp  = Inches(1.4 + row * 1.3)
    add_rect(slide, lft, tp, Inches(3.7), Inches(0.95), RGBColor(0x10,0x20,0x40))
    add_rect(slide, lft, tp, Inches(0.07), Inches(0.95), colors_pill[i])
    add_text(slide, code,   lft + Inches(0.15), tp + Inches(0.06), Inches(0.6), Inches(0.35),
             size=14, bold=True, color=colors_pill[i])
    add_text(slide, label,  lft + Inches(0.15), tp + Inches(0.44), Inches(3.4), Inches(0.4),
             size=11, color=LIGHT)

# Bottom bar
add_rect(slide, 0, SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), CYAN)

# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — SOMMARIO
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), CYAN)
add_rect(slide, 0, SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), CYAN)

add_text(slide, "Le 8 Evoluzioni Future di Metis ACE", Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
         size=28, bold=True, color=WHITE)
add_rect(slide, Inches(0.5), Inches(0.95), Inches(3), Inches(0.04), CYAN)

summaries = [
    ("E1 🔴", "Simulatore Prodotto Finanziario",    "Identifica il prodotto creditizio statisticamente più profittevole per il cliente (P-endo, P-uto, garanzie, notifiche)"),
    ("E2 🟡", "Indicatori di Pagamento per Cluster","Verifica se il comportamento di pagamento del debitore è in linea con il cluster di appartenenza"),
    ("E3 ⚪", "Teoria dei Giochi in PEF",           "Strumenti decisionali sofisticati basati su teoria dei giochi per la negoziazione delle condizioni creditizie"),
    ("E4 🔴", "Monitoraggio Predittivo (Early Warning)", "AI predittiva per l'identificazione precoce di deterioramenti su portafoglio clienti"),
    ("E5 ⚪", "Indicatori Rischio AR",               "Valutazione dell'Atipicità e Rilevanza delle operazioni per il controllo del rischio di comportamenti anomali"),
    ("E6 🟡", "Cluster Omogenei Clientela",          "Segmentazione dei clienti in cluster omogenei per individuare comportamenti anomali rispetto al gruppo"),
    ("E7 🔴", "Lettura Allegati Documentali (OCR+AI)","AI che legge, estrae e sintetizza automaticamente i contenuti degli allegati alla pratica"),
    ("E8 🟡", "Verifica Aderenza Policy Creditizie", "Analisi automatica della pratica rispetto alle policy creditizie interne della banca"),
]

for i, (code, title, desc) in enumerate(summaries):
    col = i % 2; row = i // 2
    lft = Inches(0.4 + col * 6.45)
    tp  = Inches(1.15 + row * 1.4)
    add_rect(slide, lft, tp, Inches(6.1), Inches(1.25), RGBColor(0x0D, 0x1F, 0x3C))
    add_rect(slide, lft, tp, Inches(0.06), Inches(1.25), colors_pill[i])
    add_text(slide, code,  lft+Inches(0.14), tp+Inches(0.07), Inches(1.0), Inches(0.35),
             size=11, bold=True, color=colors_pill[i])
    add_text(slide, title, lft+Inches(0.14), tp+Inches(0.35), Inches(5.8), Inches(0.35),
             size=12, bold=True, color=WHITE)
    add_text(slide, desc,  lft+Inches(0.14), tp+Inches(0.68), Inches(5.8), Inches(0.5),
             size=9,  color=LIGHT)

# ══════════════════════════════════════════════════════════════════
# SLIDES 3–10 — UNA PER EVOLUZIONE
# ══════════════════════════════════════════════════════════════════
evoluzioni = [
    {
        "code": "E1", "priority": "🔴 PRIORITÀ ALTA", "pcolor": GREEN,
        "title": "Simulatore Prodotto Finanziario",
        "img": IMAGES["e1"],
        "what": "Un motore di raccomandazione AI che, analizzato il profilo creditizio del cliente, suggerisce il prodotto finanziario statisticamente più adatto e profittevole.",
        "how": [
            "Analizza profilo cliente: settore, dimensione, storico CR, DSCR, esposizioni in essere",
            "Confronta prodotti disponibili: linea P-endo vs P-uto, con/senza garanzie, con/senza notifica",
            "Calcola profittabilità attesa e probabilità di rimborso per ogni opzione",
            "Genera raccomandazione ranked con spiegazione XAI per l'analista",
        ],
        "value": "Trasforma Metis da strumento di analisi a strumento di raccomandazione commerciale — aumenta revenue per pratica.",
    },
    {
        "code": "E2", "priority": "🟡 PRIORITÀ MEDIA", "pcolor": RGBColor(0xFF,0xC0,0x00),
        "title": "Indicatori di Pagamento per Cluster",
        "img": IMAGES["e2"],
        "what": "Sistema di benchmarking comportamentale che valuta se il pattern di pagamento del cliente è in linea con il cluster omogeneo di appartenenza.",
        "how": [
            "Segmenta il portafoglio clienti in cluster per settore, dimensione e storico",
            "Per ogni cluster calcola i KPI di comportamento di pagamento (DSO, puntualità, utilizzo)",
            "Confronta il singolo cliente con la distribuzione del cluster",
            "Genera indicatore di deviazione con severità (IN LINEA / ATTENZIONE / ANOMALO)",
        ],
        "value": "Early detection di deterioramento nascosto — il cliente appare sano ma devia dal suo cluster prima che il rating scenda.",
    },
    {
        "code": "E3", "priority": "⚪ PRIORITÀ BASSA", "pcolor": GRAY,
        "title": "Teoria dei Giochi in PEF",
        "img": IMAGES["e3"],
        "what": "Applicazione della game theory al processo negoziale del credito per ottimizzare le condizioni offerte dalla banca in funzione delle possibili contromosses del cliente.",
        "how": [
            "Modella la negoziazione come gioco a 2 giocatori (Banca vs Cliente) con payoff definiti",
            "Calcola le strategie dominanti e l'equilibrio di Nash in base al profilo del cliente",
            "Suggerisce la struttura ottimale dell'offerta (tasso, garanzie, durata, covenants)",
            "Simula scenari controfattuali: 'cosa succede se il cliente rinegozia?'",
        ],
        "value": "Strumento avanzato per analisti senior su pratiche complesse >1M€. Differenziatore unico sul mercato bancario italiano.",
    },
    {
        "code": "E4", "priority": "🔴 PRIORITÀ ALTA", "pcolor": GREEN,
        "title": "Monitoraggio Predittivo — Early Warning AI",
        "img": IMAGES["e4"],
        "what": "Sistema di AI predittiva che monitora in continuo il portafoglio crediti e identifica precocemente i clienti a rischio di deterioramento nei prossimi 30–90 giorni.",
        "how": [
            "Feed continuo di dati: CR mensile, utilizzo linee, pagamenti, news web (M2)",
            "Modello ML (XGBoost + LSTM) addestrato su default storici",
            "Score di probabilità di default a 30/60/90 giorni con SHAP explanation",
            "Alert automatici all'analista responsabile con dashboard aggregata portfolio",
        ],
        "value": "Anticipa le crisi di 60–90 giorni rispetto al rilevamento tradizionale. Riduce le perdite su crediti (LGD). Massimo ritorno sull'investimento.",
    },
    {
        "code": "E5", "priority": "⚪ PRIORITÀ BASSA", "pcolor": GRAY,
        "title": "Indicatori di Rischio AR (Atipicità/Rilevanza)",
        "img": IMAGES["e5"],
        "what": "Modulo per la valutazione del rischio di comportamenti atipici e operazioni di rilevanza anomala, a supporto dei processi antiriciclaggio e di credit risk monitoring.",
        "how": [
            "Analizza pattern transazionali: frequenza, importi, controparti, orari",
            "Calcola indice di Atipicità (deviazione dalla norma comportamentale del cluster)",
            "Calcola indice di Rilevanza (impatto dell'operazione sul profilo di rischio)",
            "Genera segnalazione strutturata per il compliance officer",
        ],
        "value": "Supporto al processo AML (Anti Money Laundering) integrato nel flusso di underwriting. Riduce il rischio operativo e reputazionale.",
    },
    {
        "code": "E6", "priority": "🟡 PRIORITÀ MEDIA", "pcolor": RGBColor(0xFF,0xC0,0x00),
        "title": "Cluster Omogenei per Early Warning Comportamentale",
        "img": IMAGES["e6"],
        "what": "Segmentazione non supervisionata del portafoglio clienti in cluster omogenei per caratteristiche economiche e comportamentali, con monitoraggio delle deviazioni.",
        "how": [
            "Algoritmo di clustering (K-Means / DBSCAN) su KPI di bilancio + CR + settore",
            "Aggiornamento mensile dei cluster all'arrivo dei nuovi dati CR",
            "Monitoraggio continuo: ogni cliente viene comparato con l'evoluzione del suo cluster",
            "Alert quando un cliente 'migra' verso un cluster a rischio più elevato",
        ],
        "value": "Base scientifica per l'early warning. I cluster definiscono il 'normale' — le deviazioni diventano automaticamente segnali di allarme.",
    },
    {
        "code": "E7", "priority": "🔴 PRIORITÀ ALTA", "pcolor": GREEN,
        "title": "Lettura e Sintesi Allegati Documentali (OCR + AI)",
        "img": IMAGES["e7"],
        "what": "Modulo che legge automaticamente i documenti PDF allegati alla pratica (bilanci, visure, contratti, perizie) e ne estrae gli elementi critici evidenziandoli all'analista.",
        "how": [
            "OCR ad alta precisione su PDF scansionati (LlamaParse / Tesseract)",
            "Estrazione strutturata di dati chiave: date, importi, soggetti, clausole",
            "AI sintesi: 'Questo contratto di leasing ha una clausola di accelerazione al downgrade'",
            "Highlight automatico degli elementi che richiedono approfondimento",
        ],
        "value": "Elimina ore di lettura manuale di documenti. Riduce il rischio di tralasciare clausole critiche. Direttamente integrabile nella pipeline OCR già presente.",
    },
    {
        "code": "E8", "priority": "🟡 PRIORITÀ MEDIA", "pcolor": RGBColor(0xFF,0xC0,0x00),
        "title": "Verifica Aderenza alle Policy Creditizie",
        "img": IMAGES["e8"],
        "what": "Modulo che verifica automaticamente se il contenuto e le conclusioni della pratica sono conformi alle policy creditizie interne della banca.",
        "how": [
            "Caricamento del regolamento crediti interno come knowledge base (RAG)",
            "Analisi della pratica rispetto alle regole: limite di concentrazione, settori esclusi, LTV, covenant",
            "Checklist automatica con esito per ogni regola (✅ conforme / ❌ violazione / ⚠️ attenzione)",
            "Score di compliance complessivo con evidenza delle violazioni per il credit officer",
        ],
        "value": "Standardizza la qualità delle pratiche. Riduce il rischio operativo da errori umani su regole complesse. Si integra con l'Audit Trail esistente.",
    },
]

for evo in evoluzioni:
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide)

    # Top accent bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.05), evo["pcolor"])

    # Left panel — testo
    add_rect(slide, 0, 0, Inches(5.6), SLIDE_H, RGBColor(0x05, 0x12, 0x2A))

    # Code badge
    add_rect(slide, Inches(0.3), Inches(0.2), Inches(0.65), Inches(0.45), evo["pcolor"])
    add_text(slide, evo["code"], Inches(0.3), Inches(0.2), Inches(0.65), Inches(0.45),
             size=16, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

    # Priority
    add_text(slide, evo["priority"], Inches(1.1), Inches(0.22), Inches(4.0), Inches(0.4),
             size=11, bold=True, color=evo["pcolor"])

    # Title
    add_text(slide, evo["title"], Inches(0.3), Inches(0.75), Inches(5.0), Inches(0.85),
             size=21, bold=True, color=WHITE)

    # Divider
    add_rect(slide, Inches(0.3), Inches(1.6), Inches(2.5), Inches(0.035), evo["pcolor"])

    # Cos'è
    add_text(slide, "COS'È", Inches(0.3), Inches(1.75), Inches(5.0), Inches(0.35),
             size=9, bold=True, color=evo["pcolor"])
    add_text(slide, evo["what"], Inches(0.3), Inches(2.05), Inches(5.0), Inches(1.0),
             size=10.5, color=LIGHT)

    # Come funziona
    add_text(slide, "COME FUNZIONA", Inches(0.3), Inches(3.1), Inches(5.0), Inches(0.35),
             size=9, bold=True, color=evo["pcolor"])
    for j, step in enumerate(evo["how"]):
        add_text(slide, f"{'①②③④'[j]}  {step}",
                 Inches(0.3), Inches(3.4 + j*0.68), Inches(5.1), Inches(0.6),
                 size=10, color=WHITE)

    # Valore
    add_rect(slide, Inches(0.3), Inches(6.55), Inches(5.0), Inches(0.7), RGBColor(0x0A,0x28,0x4A))
    add_text(slide, "💡 " + evo["value"],
             Inches(0.4), Inches(6.6), Inches(4.9), Inches(0.65),
             size=9.5, italic=True, color=GOLD)

    # Right panel — image
    if os.path.exists(evo["img"]):
        slide.shapes.add_picture(
            evo["img"],
            Inches(5.75), Inches(0.55),
            Inches(7.3), Inches(6.7)
        )

    # Bottom bar
    add_rect(slide, 0, SLIDE_H - Inches(0.05), SLIDE_W, Inches(0.05), evo["pcolor"])

# ══════════════════════════════════════════════════════════════════
# SLIDE FINALE — ROADMAP TIMELINE
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.05), CYAN)
add_rect(slide, 0, SLIDE_H - Inches(0.05), SLIDE_W, Inches(0.05), CYAN)

add_text(slide, "Roadmap di Sviluppo", Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
         size=28, bold=True, color=WHITE)
add_rect(slide, Inches(0.5), Inches(0.88), Inches(2.8), Inches(0.04), CYAN)

# Timeline horizontal line
add_rect(slide, Inches(0.5), Inches(4.3), Inches(12.2), Inches(0.06), RGBColor(0x1E,0x40,0x80))

phases = [
    {"label": "FASE 1\nQ3 2026", "items": ["E7 — OCR Allegati", "E4 — Early Warning"], "color": GREEN, "x": 0.5},
    {"label": "FASE 2\nQ4 2026", "items": ["E6 — Cluster Clientela", "E1 — Simulatore Prodotto"], "color": CYAN, "x": 3.8},
    {"label": "FASE 3\nQ1 2027", "items": ["E8 — Verifica Policy", "E2 — Cluster Pagamento"], "color": RGBColor(0xFF,0xC0,0x00), "x": 7.1},
    {"label": "FASE 4\nQ2 2027", "items": ["E5 — Rischio AR", "E3 — Teoria Giochi"], "color": GRAY, "x": 10.4},
]

for ph in phases:
    lft = Inches(ph["x"])
    # dot on timeline
    add_rect(slide, lft + Inches(1.1), Inches(4.18), Inches(0.3), Inches(0.3), ph["color"])
    # phase label
    add_text(slide, ph["label"], lft, Inches(4.55), Inches(2.8), Inches(0.8),
             size=11, bold=True, color=ph["color"], align=PP_ALIGN.CENTER)
    # items
    for k, item in enumerate(ph["items"]):
        add_rect(slide, lft, Inches(2.6 + k*0.85), Inches(2.8), Inches(0.7), RGBColor(0x0D,0x1F,0x3C))
        add_rect(slide, lft, Inches(2.6 + k*0.85), Inches(0.05), Inches(0.7), ph["color"])
        add_text(slide, item, lft+Inches(0.12), Inches(2.65 + k*0.85), Inches(2.6), Inches(0.55),
                 size=10, color=WHITE)

add_text(slide, "Metis ACE — Roadmap 2026–2027  |  Soggetto a revisione trimestrale",
         Inches(0.5), Inches(6.9), Inches(12), Inches(0.4),
         size=9, color=GRAY, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────
out = r"c:\Users\GaetanoPecorella\OneDrive - Altermaind\Desktop\METIS\Metis_Roadmap_Evoluzioni.pptx"
prs.save(out)
import shutil
shutil.copy(out, r"C:\Users\GaetanoPecorella\Downloads\Metis_Roadmap_Evoluzioni.pptx")
print("Salvato:", out)
print("Copiato in Downloads.")
