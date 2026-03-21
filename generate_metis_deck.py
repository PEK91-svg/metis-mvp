import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# McKinsey Colors
PRIMARY_BLUE = RGBColor(0, 41, 96)
SECONDARY_BLUE = RGBColor(0, 101, 189)
LIGHT_BLUE = RGBColor(201, 240, 255)
YELLOW = RGBColor(255, 219, 0)
WHITE = RGBColor(255, 255, 255)

def add_title(slide, text):
    title_shape = slide.shapes.title
    title_shape.text = text
    title_shape.text_frame.paragraphs[0].font.size = Pt(18)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    title_shape.text_frame.word_wrap = True

def add_content_box(slide, left, top, width, height, text, bg_color=None, font_color=None, font_size=11, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    if font_color:
        p.font.color.rgb = font_color
    p.alignment = align
    
    if bg_color:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    return txBox

prs = Presentation()
# Set 16:9 
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

blank_slide_layout = prs.slide_layouts[6]
title_only_layout = prs.slide_layouts[5]

# --- Slide 1: Cover ---
slide1 = prs.slides.add_slide(blank_slide_layout)
bg = slide1.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625)) 
bg.fill.solid()
bg.fill.fore_color.rgb = PRIMARY_BLUE
bg.line.fill.background()

add_content_box(slide1, 1, 2, 8, 1, "Metis - Analisi Complementare Evoluta", None, WHITE, 28, PP_ALIGN.CENTER)
add_content_box(slide1, 1, 3, 8, 0.5, "L'Intelligenza Aumentata (XAI) per l'innovazione e la conformità nel Credit Underwriting Bancario.", None, WHITE, 14, PP_ALIGN.CENTER)

# --- Slide 2: Executive Summary ---
slide2 = prs.slides.add_slide(title_only_layout)
add_title(slide2, "L'Automazione \"Glass Box\" di Metis riduce i tempi istruttori garantendo compliance normativa EU")
# 3 boxes
add_content_box(slide2, 0.5, 1.5, 2.8, 3, "IL PROBLEMA\n\nIstruttorie PEF lente, data-entry manuale tra bilanci e Centrale Rischi. Rischio di bias cognitivi.", LIGHT_BLUE, PRIMARY_BLUE, 11)
add_content_box(slide2, 3.6, 1.5, 2.8, 3, "LA SOLUZIONE METIS\n\nPiattaforma SaaS con 8 moduli di AI spiegabile (XAI). Si affianca all'analista offrendo insights trasparenti.", PRIMARY_BLUE, WHITE, 11)
add_content_box(slide2, 6.7, 1.5, 2.8, 3, "IL VALORE CREATO\n\nIncremento AuROC: +0.6 - 1.9%\nRiduzione Tempi: -70%\nCompliance EU AI Act: 100%", SECONDARY_BLUE, WHITE, 11)

# --- Slide 3: Il problema ---
slide3 = prs.slides.add_slide(title_only_layout)
add_title(slide3, "Il processo PEF attuale è frammentato, lento e soggetto a crescenti rischi operativi")
add_content_box(slide3, 0.5, 1.5, 4.5, 3.5, "DIAGRAMMA DEL COLLO DI BOTTIGLIA\n\n- Raccolta visura Centrale Rischi (PDF)\n- Estrazione dati da Bilancio (Excel)\n- Analisi web sulle news (Google)\n- Stesura report manuale (Word)\n\nL'analista perde il 70% del tempo nella raccolta dati, e solo il 30% nell'analisi del credito.", None, PRIMARY_BLUE, 11)
add_content_box(slide3, 5.5, 1.5, 4, 1.5, "RISCHIO QUALITATIVO (Black Box Risk)\nDifficoltà metodologica per modelli non standard.", YELLOW, PRIMARY_BLUE, 11)
add_content_box(slide3, 5.5, 3.2, 4, 1.5, "RISCHIO TEMPI E COSTI (Data Fragmentation)\nPerdita di competitività commerciale sulle tempistiche di delibera.", YELLOW, PRIMARY_BLUE, 11)

# --- Slide 4: EU AI Act ---
slide4 = prs.slides.add_slide(title_only_layout)
add_title(slide4, "L'EU AI Act classifica i modelli di Credit Scoring come 'Alto Rischio', vietando di fatto le scatole nere")
add_content_box(slide4, 0.5, 1.5, 9, 1, "A partire da Agosto 2026, l'uso dell'Intelligenza Artificiale nel credito bancario sarà soggetto ai requisiti severi di 'Explainable AI' (XAI).", LIGHT_BLUE, PRIMARY_BLUE, 14, PP_ALIGN.CENTER)
add_content_box(slide4, 0.5, 3, 2.8, 2, "TRASPARENZA\n\nI modelli devono spiegare testualmente da quali fattori è scaturita un'approvazione o un diniego.", PRIMARY_BLUE, WHITE)
add_content_box(slide4, 3.6, 3, 2.8, 2, "HUMAN-IN-THE-LOOP\n\nL'IA raccomanda, l'uomo delibera. Necessità di una UX che lasci l'analista al centro del controllo.", PRIMARY_BLUE, WHITE)
add_content_box(slide4, 6.7, 3, 2.8, 2, "AUDITABILITY\n\nTutte le decisioni dell'IA devono esplicitare in chiaro la fonte (es. 'Articolo Sole24Ore' o 'Pagina Bilancio').", PRIMARY_BLUE, WHITE)

# --- Slide 5: La Soluzione ---
slide5 = prs.slides.add_slide(title_only_layout)
add_title(slide5, "Metis risolve il dilemma normativo offrendo Intelligenza Aumentata ('Glass Box') nativa")
add_content_box(slide5, 0.5, 1.5, 4, 3, "SCENARIO ATTUALE\n\n- Tempi decisonali lunghi\n- Modelli legacy puramente quantitativi\n- Analisi qualitativa puramente umana", YELLOW, PRIMARY_BLUE, 12, PP_ALIGN.CENTER)
add_content_box(slide5, 4.6, 2.5, 0.8, 1, "=>", None, PRIMARY_BLUE, 24, PP_ALIGN.CENTER)
add_content_box(slide5, 5.5, 1.5, 4, 3, "SCENARIO METIS\n\n- Copilot basato su GenAI (XAI)\n- 8 Moduli che incrociano testo, andamenti e bilanci\n- Suggerimenti documentali citati e validabili al 100% dall'analista.", PRIMARY_BLUE, WHITE, 12, PP_ALIGN.CENTER)

# --- Slide 6: Architettura 8 Moduli ---
slide6 = prs.slides.add_slide(title_only_layout)
add_title(slide6, "Un'architettura basata su 8 Moduli Core automatizza completamente l'output narrativo del fido")
add_content_box(slide6, 0.5, 1.2, 9, 0.5, "MODULI DI ANALISI DELLA DOCUMENTAZIONE E SCORING", PRIMARY_BLUE, WHITE, 12, PP_ALIGN.CENTER)

cols = ["1. Riepilogo Storico", "2. Web Sentiment", "3. Analisi KPI", "4. Benchmark ISTAT", 
        "5. Pattern CR (12 m)", "6. Cross-check dati", "7. Forecast (DSCR)", "8. Analisi SWOT"]
for i, name in enumerate(cols):
    row = i // 4
    col = i % 4
    add_content_box(slide6, 0.5 + (col*2.3), 1.9 + (row*1.5), 2.1, 1.2, name, LIGHT_BLUE, PRIMARY_BLUE, 12, PP_ALIGN.CENTER)

# --- Slide 7: Superiorità ---
slide7 = prs.slides.add_slide(title_only_layout)
add_title(slide7, "La superiorità analitica: incremento della capacità discriminante e abbattimento dei tempi istruttori")
add_content_box(slide7, 1, 2, 3.5, 2.5, "AUMENTO CAPACITÀ DISCRIMINANTE\n(Backtesting)\n\n+0,6% -> +1,9% (AuROC migliorato)", SECONDARY_BLUE, WHITE, 14, PP_ALIGN.CENTER)
add_content_box(slide7, 5.5, 2, 3.5, 2.5, "RIDUZIONE TEMPI (EFFICIENZA)\n\nTempi di redazione della singola PEF ridotti mediamente del 70%.", PRIMARY_BLUE, WHITE, 14, PP_ALIGN.CENTER)

# --- Slide 8: Competitor ---
slide8 = prs.slides.add_slide(title_only_layout)
add_title(slide8, "Rispetto ad alternative legacy e start-up, Metis è l'unica focalizzata sul workflow 'Glass Box'")
# Very basic representation of a table using text boxes
add_content_box(slide8, 0.5, 1.5, 2.8, 3.5, "PLAYER LEGACY\n(CRIF, Experian)\n\n• Score numerico: Sì\n• Focus normativo: Alto\n• Modello Innovativo NLP: Basso\n• Generazione testo per PEF: Basso", None, PRIMARY_BLUE, 12)
add_content_box(slide8, 3.6, 1.5, 2.8, 3.5, "STARTUP AI / DECISION ENGINE\n(Actico, Axyon)\n\n• Score numerico: Sì\n• Focus automazione drastica: Alto\n• Modello Spiegabile 'Narrativo': Medio-Basso", None, PRIMARY_BLUE, 12)
add_content_box(slide8, 6.7, 1.5, 2.8, 3.5, "METIS (XAI)\n\n• Score numerico + Testo Narrativo: Sì\n• Focus Normativo Alto Rischio (EU Act): Altissimo\n• Approccio Human-In-The-Loop: Nativo", LIGHT_BLUE, PRIMARY_BLUE, 12)

# --- Slide 9: Business Model ---
slide9 = prs.slides.add_slide(title_only_layout)
add_title(slide9, "Il Business Model ibrido bilancia costi prevedibili e adozione scalabile Pay-Per-Use")
add_content_box(slide9, 0.5, 1.5, 2.8, 3, "1. BASE PLATFORM FEE\n\nLicenza SaaS annuale. Copre oneri di conformità, update modelli genAI e cloud hosting sicuro (Mantenimento flat).", PRIMARY_BLUE, WHITE, 12)
add_content_box(slide9, 3.6, 1.5, 2.8, 3, "2. CONSUMPTION TIER (Pay-Per-PEF)\n\nFee variabile in base al volume delle pratiche. Assicura ROI sempre maggiore: la banca paga solo se Metis abbatte il carico di lavoro.", YELLOW, PRIMARY_BLUE, 12)
add_content_box(slide9, 6.7, 1.5, 2.8, 3, "3. ENTERPRISE ADD-ON\n\nServizi professionali ad hoc. Addestramento fine-tuned sui dati storici delle PEF interne della singola banca per personalizzazione UX.", SECONDARY_BLUE, WHITE, 12)

# --- Slide 10: Conclusioni ---
slide10 = prs.slides.add_slide(title_only_layout)
add_title(slide10, "Sintesi strategica: Call to Action per il lancio della Fase Pilota sulla piattaforma")
add_content_box(slide10, 0.5, 1.5, 4.5, 3.5, "I 3 TAKEAWAY DEL PROGETTO:\n\n1. Il mercato bancario ha urgenze di efficienza ed è sotto pressione normativa dal 2026.\n\n2. Metis risponde nativamente con il framework Explainable AI.\n\n3. L'analista viene 'aumentato', minimizzando la resistenza sindacale/culturale interna.", PRIMARY_BLUE, WHITE, 14)
add_content_box(slide10, 5.5, 1.5, 4, 3.5, "PROSSIMI PASSI (TIMELINE PoC)\n\nMese 1: Setup ambiente e ingestion dati di test CR/Bilanci\n\nMese 2: Back-testing cieco su PEF storiche\n\nMese 3: Valutazione KPI e rollout commerciale", LIGHT_BLUE, PRIMARY_BLUE, 14)

prs.save("Metis_Pitch_Deck.pptx")
print("Presentation generated successfully at Metis_Pitch_Deck.pptx")
